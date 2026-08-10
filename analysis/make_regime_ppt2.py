"""레짐 모델 탐색과정 PPT (3장, 그래프+수치 중심) — 박사님 보고용. 문제의식 슬라이드 제외."""
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

CH = Path("analysis/ppt_charts"); CH.mkdir(exist_ok=True)
GREEN = "#1B6E3C"; RED = "#B12A2A"; BLUE = "#2E5B9A"; GRAY = "#777777"; ORANGE = "#C77A14"

# ---------- Chart 1: seed robustness (why FX adopted) ----------
seeds = ["0", "1", "7", "42", "123"]
hsmm = [-0.93, -1.75, -2.18, 1.75, -1.91]   # breadth-only HSMM
hsmm_fx = [1.84, 1.53, 1.53, 1.36, 1.53]    # +FX
x = np.arange(len(seeds)); w = 0.38
fig, ax = plt.subplots(figsize=(6.6, 4.0))
ax.bar(x - w/2, hsmm, w, label="breadth only", color=RED, alpha=.85)
ax.bar(x + w/2, hsmm_fx, w, label="+ FX", color=GREEN, alpha=.9)
ax.axhline(0, color="k", lw=.8)
ax.set_xticks(x); ax.set_xticklabels(seeds); ax.set_xlabel("random seed")
ax.set_ylabel("Bull/Bear gap (%p)")
ax.set_title("Why FX was adopted: stabilizes discrimination\n(breadth-only flips sign by seed; +FX all +1.5p)", fontsize=10)
ax.legend(fontsize=9); fig.tight_layout(); fig.savefig(CH/"c1_seed.png", dpi=150); plt.close()

# ---------- Chart 2: config summary (gap mean ± std) ----------
cfgs = ["base\n(+FX)", "+US10Y", "+DXY", "+News\n(T+S)", "3-state", "4-state"]
mu = [1.56, 2.32, 1.77, -0.80, -0.58, 0.43]
sd = [0.16, 0.12, 0.00, 3.29, 1.82, 1.90]
note = ["CHOSEN", "no FCF gain", "weak", "unstable", "unstable", "unstable"]
cols = [GREEN, ORANGE, BLUE, RED, RED, RED]
fig, ax = plt.subplots(figsize=(6.6, 4.0))
xx = np.arange(len(cfgs))
ax.bar(xx, mu, yerr=sd, color=cols, alpha=.85, capsize=4)
ax.axhline(0, color="k", lw=.8)
for i, (m, s, t) in enumerate(zip(mu, sd, note)):
    ax.text(i, m + (s if m >= 0 else -s) + 0.15*np.sign(m or 1), t, ha="center", fontsize=8, color=cols[i])
ax.set_xticks(xx); ax.set_xticklabels(cfgs, fontsize=8.5)
ax.set_ylabel("Bull/Bear gap, 5-seed μ ± σ (%p)")
ax.set_title("Feature/structure exploration: only base(+FX) is\nstable & positive; others blow up σ or go negative", fontsize=10)
fig.tight_layout(); fig.savefig(CH/"c2_config.png", dpi=150); plt.close()

# ---------- Chart 3: vs ai_v2 ----------
fig, axs = plt.subplots(1, 2, figsize=(7.2, 3.8))
# (a) regime quality
m1 = ["Recall\n(/7)", "Gap (%p)"]; hs = [7, 1.53]; ai = [6, 0.23]
xx = np.arange(2); w = .36
axs[0].bar(xx-w/2, hs, w, label="hsmm+FX", color=GREEN); axs[0].bar(xx+w/2, ai, w, label="ai_v2", color=GRAY)
axs[0].set_xticks(xx); axs[0].set_xticklabels(m1, fontsize=8.5); axs[0].set_title("Regime quality (hsmm wins)", fontsize=9.5); axs[0].legend(fontsize=8)
# (b) FCF
m2 = ["Cum.ret (%)", "|MDD| (%)"]; hs2 = [353, 38.1]; ai2 = [406, 40.6]
axs[1].bar(xx-w/2, hs2, w, label="hsmm+FX", color=GREEN); axs[1].bar(xx+w/2, ai2, w, label="ai_v2", color=GRAY)
axs[1].set_xticks(xx); axs[1].set_xticklabels(m2, fontsize=8.5); axs[1].set_title("FCF (ai_v2 higher return)", fontsize=9.5); axs[1].legend(fontsize=8)
fig.suptitle("hsmm+FX vs ai_v2 (overlap 2018~2026)", fontsize=11); fig.tight_layout(); fig.savefig(CH/"c3_aiv2.png", dpi=150); plt.close()

# ================= PPT =================
NAVY = RGBColor(0x1F,0x33,0x55); BLUEc = RGBColor(0x2E,0x5B,0x9A); GR = RGBColor(0x55,0x55,0x55)
REDc = RGBColor(0xB1,0x2A,0x2A); GRNc = RGBColor(0x1B,0x6E,0x3C)
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)


def slide(title, sub):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tf = s.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.95)).text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = NAVY
    p2 = tf.add_paragraph(); p2.text = sub; p2.font.size = Pt(12.5); p2.font.italic = True; p2.font.color.rgb = BLUEc
    return s


def body(s, items, left, top, width, size=12.5):
    tf = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.8)).text_frame; tf.word_wrap = True
    for i, (lvl, txt, col, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.level = lvl; p.font.size = Pt(size - lvl); p.font.bold = bold; p.font.color.rgb = col; p.space_after = Pt(2)


# Slide 1 — 피처 탐색 + Chart2
s = slide("레짐 모델 탐색 ① — 피처 추가 실험 (변수도 추가해봄)",
          "base = HSMM(2-state)+환율. 추가 피처가 분별력/안정성을 개선하는지 5-seed로 검증")
body(s, [
 (0,"■ 구성별 Bull/Bear 수익격차 (5-seed μ±σ, overlap)",NAVY,True),
 (1,"base(+환율)  +1.56 ±0.16   → 안정·기준",GRNc,False),
 (1,"+US10Y(금리) +2.32 ±0.12   → 분별력 최고지만 FCF 수익 무익 (격차≠수익)",REDc,False),
 (1,"+DXY(달러)   +1.77 ±0.00   → 약효과",GR,False),
 (1,"+News(T+S)   −0.80 ±3.29   → robustness 파괴·2023 과잉방어",REDc,False),
 (1,"구리/금       독립정보 77%이나 bear 예측 증분 AUC −0.007·후행 → 무익",REDc,False),
 (0,"■ 반복 확인된 교훈",NAVY,True),
 (1,"분별력(격차) ↑ ≠ FCF 수익 ↑   /   독립 ≠ 유용",REDc,True),
 (1,"→ 피처 추가만으로는 base 못 이김",GR,False),
], 0.45, 1.45, 6.2)
s.shapes.add_picture(str(CH/"c2_config.png"), Inches(6.9), Inches(2.0), width=Inches(6.1))

# Slide 2 — 환율 robustness + 구조/readout + Chart1
s = slide("레짐 모델 탐색 ② — 핵심 발견(환율) · 구조/예측방식 시도",
          "왜 2-state HSMM+환율을 최종 선택했나")
body(s, [
 (0,"■ 결정적 발견 — seed 취약성",NAVY,True),
 (1,"breadth만: 격차가 seed마다 부호 뒤집힘 (−2.18 ~ +1.75) = 분별력이 운",REDc,False),
 (1,"+환율: 5-seed 전부 +1.4~+1.8 (옆 그래프) → robustness 확보",GRNc,True),
 (1,"환율 = 예측력 아니라 '자본흐름 독립축'으로 상태 식별 안정화",GR,False),
 (0,"■ 구조 시도 (상태 수)",NAVY,True),
 (1,"3-state −0.58±1.82 / 4-state +0.43±1.90 → 상태는 분화되나 격차 불안정 → 기각",REDc,False),
 (0,"■ 예측방식 시도",NAVY,True),
 (1,"전이행렬 1-step: Lead +0.7 vs persistence +0.9 → 개선 없음 → 기각",REDc,False),
 (0,"■ 최종 선택",NAVY,True),
 (1,"Full-Cov HSMM(2-state)+환율  ·  피처 6개(breadth5+환율)  ·  5-seed consensus",GRNc,True),
], 0.45, 1.45, 6.2)
s.shapes.add_picture(str(CH/"c1_seed.png"), Inches(6.9), Inches(2.2), width=Inches(6.1))

# Slide 3 — 결과 vs ai_v2 + Chart3
s = slide("레짐 모델 ③ — 결과 vs ai_v2 · 한계 · 다음",
          "overlap 2018~2026")
body(s, [
 (0,"■ 성능",NAVY,True),
 (1,"Recall 7/7 vs 6/7 · 격차 +1.53 vs +0.23 · MDD −38% vs −41% → 위기탐지·분별력·방어 우위",GRNc,False),
 (1,"FCF 누적 353% vs 406% → 절대수익은 ai_v2 우위",REDc,False),
 (0,"■ ai_v2 수익이 높은 이유 (attribution)",NAVY,True),
 (1,"① 2022 긴축: ai 지속 Bear(매크로) vs hsmm breadth 깜빡임 (−7.4%p)",GR,False),
 (1,"② 회복장: hsmm 바닥 손절 139건 vs ai Bull 유지 (−5.2%p)",GR,False),
 (1,"③ ai Bear월·Whipsaw 적음 → 종목/분별력 아니라 '정보·타이밍'",GR,False),
 (0,"■ 한계 & 다음",NAVY,True),
 (1,"ai 우위 = 뉴스/매크로 '정보' → HMM 피처 주입 모두 실패",REDc,True),
 (1,"다음: Risk overlay(강세전략 유지+Bear 익스포저 축소) + ai 하이브리드 게이트",BLUEc,True),
], 0.45, 1.45, 6.4)
s.shapes.add_picture(str(CH/"c3_aiv2.png"), Inches(7.1), Inches(2.4), width=Inches(5.9))

out = Path("레짐예측모델_탐색정리.pptx"); prs.save(out)
print(f"저장: {out.resolve()} (3장) + 차트 3개 {CH.resolve()}")
