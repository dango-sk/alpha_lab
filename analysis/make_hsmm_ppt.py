# -*- coding: utf-8 -*-
"""
analysis/make_hsmm_ppt.py
HSMM 레짐 2차 보고 PPT 생성 (박사님 피드백 반영 + Risk Overlay 검증).
차트 2개(MDD 비교 / 익스포저 중복) + 8슬라이드. 출력: HSMM_레짐_2차_보고.pptx
사용: .venv/bin/python analysis/make_hsmm_ppt.py
"""
import numpy as np, pandas as pd
from pathlib import Path
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
plt.rcParams["font.family"] = "AppleGothic"; plt.rcParams["axes.unicode_minus"] = False
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = Path(__file__).parent.parent
A = Path(__file__).parent
ASSET = A / "ppt_assets"; ASSET.mkdir(exist_ok=True)
OUT = BASE / "HSMM_레짐_2차_보고.pptx"
FONT = "Apple SD Gothic Neo"

NAVY = RGBColor(0x1F, 0x38, 0x64); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x26, 0x2A, 0x33); GRAY = RGBColor(0x59, 0x5F, 0x6B)
BLUE = RGBColor(0x2F, 0x54, 0x96); GREEN = RGBColor(0x2E, 0x74, 0x4F)
RED = RGBColor(0xC0, 0x3A, 0x3A); ALT = RGBColor(0xF2, 0xF4, 0xF7)
HLG = RGBColor(0xE2, 0xEF, 0xDA)

# ───────────── 차트 ─────────────
df = pd.read_csv(A / "fcf_overlay_series.csv")
corr = np.corrcoef(df.expA, df.expP)[0, 1]

# chart1: MDD 비교
fig, ax = plt.subplots(figsize=(5.2, 3.1), dpi=200)
names = ["FCF불\n단독(BM)", "A:20일\n실현", "B:60일\n하방", "pbear만\n(vol無)"]
mdd = [-43.7, -34.5, -32.5, -33.7]
cols = [RED.rgb if False else "#C03A3A", "#2E744F", "#2E744F", "#2E744F"]
bars = ax.bar(names, mdd, color=cols, width=0.62)
for b, v in zip(bars, mdd):
    ax.text(b.get_x() + b.get_width() / 2, v - 1.5, f"{v:.1f}%", ha="center", va="top",
            color="white", fontsize=10, fontweight="bold")
ax.axhline(0, color="#888", lw=0.8)
ax.set_ylabel("최대낙폭 MDD (%)"); ax.set_ylim(-50, 2)
ax.set_title("레짐 오버레이 MDD 방어 (BM 대비 +12~13%p, P≥0.97 유의)", fontsize=10.5, pad=8)
ax.spines[["top", "right"]].set_visible(False)
fig.tight_layout(); fig.savefig(ASSET / "mdd.png", bbox_inches="tight"); plt.close(fig)

# chart2: 익스포저 중복
fig, ax = plt.subplots(figsize=(5.6, 3.1), dpi=200)
x = np.arange(len(df))
ax.plot(x, df.expP, color="#1F3864", lw=1.6, label="pbear만 (vol 無)")
ax.plot(x, df.expA, color="#2E744F", lw=1.3, ls="--", label="A: 20일 실현 vol-targeting")
yrs = df.ym.str.slice(0, 4).values
ticks = [i for i in range(len(df)) if i == 0 or yrs[i] != yrs[i - 1]]
ax.set_xticks(ticks); ax.set_xticklabels([yrs[i] for i in ticks], fontsize=8)
ax.set_ylabel("주식 익스포저"); ax.set_ylim(0.1, 1.05)
ax.set_title(f"익스포저 경로 거의 동일 — corr = {corr:.3f}\n(vol-targeting이 레짐 신호와 중복)", fontsize=10.5, pad=8)
ax.legend(fontsize=8.5, loc="lower right"); ax.spines[["top", "right"]].set_visible(False)
fig.tight_layout(); fig.savefig(ASSET / "exp.png", bbox_inches="tight"); plt.close(fig)

# ───────────── PPT 헬퍼 ─────────────
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    return prs.slides.add_slide(BLANK)


def tb(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    b = s.shapes.add_textbox(l, t, w, h); tf = b.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05); tf.margin_top = tf.margin_bottom = Inches(0.02)
    return b, tf


def run(p, text, size, color=INK, bold=False, italic=False):
    r = p.add_run(); r.text = text; f = r.font
    f.size = Pt(size); f.name = FONT; f.bold = bold; f.italic = italic; f.color.rgb = color
    return r


def header(s, num, title, subtitle=None):
    _, tf = tb(s, Inches(0.55), Inches(0.28), Inches(12.3), Inches(1.0))
    p = tf.paragraphs[0]; run(p, f"{num} ", 30, BLUE, True); run(p, title, 30, NAVY, True)
    if subtitle:
        p2 = tf.add_paragraph(); run(p2, subtitle, 13.5, GRAY, italic=True)
    ln = s.shapes.add_shape(1, Inches(0.55), Inches(1.28), Inches(12.25), Pt(2.2))
    ln.fill.solid(); ln.fill.fore_color.rgb = NAVY; ln.line.fill.background()


def bullets(s, items, l, t, w, h, size=13.5, gap=6):
    _, tf = tb(s, l, t, w, h)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(it, tuple):
            mark, txt, col, bold = it
            run(p, mark + "  ", size, col, True); run(p, txt, size, INK, bold)
        else:
            run(p, "■  ", size, NAVY, True); run(p, it, size, INK)
    return tf


def table(s, rows, l, t, w, colw, hl=None, header_h=0.42, row_h=0.44, fs=11.5):
    nr, nc = len(rows), len(rows[0])
    gr = s.shapes.add_table(nr, nc, l, t, w, Inches(header_h + row_h * (nr - 1))).table
    gr.first_row = False; gr.horz_banding = False
    tot = sum(colw)
    for j, cw in enumerate(colw):
        gr.columns[j].width = Emu(int(w * cw / tot))
    gr.rows[0].height = Inches(header_h)
    for i in range(1, nr):
        gr.rows[i].height = Inches(row_h)
    for i, rowd in enumerate(rows):
        for j, val in enumerate(rowd):
            c = gr.cell(i, j); c.margin_left = c.margin_right = Inches(0.08)
            c.margin_top = c.margin_bottom = Inches(0.03); c.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = NAVY
            elif hl is not None and i == hl:
                c.fill.solid(); c.fill.fore_color.rgb = HLG
            else:
                c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else ALT
            tfc = c.text_frame; tfc.word_wrap = True
            p = tfc.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            col = WHITE if i == 0 else (GREEN if (hl and i == hl) else INK)
            run(p, str(val), fs, col, bold=(i == 0 or (hl and i == hl)))
    return gr


def conclusion(s, parts, t, l=Inches(0.55), w=Inches(12.25), size=13.5):
    _, tf = tb(s, l, t, w, Inches(1.1))
    for i, (txt, col, bold) in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(4)
        run(p, "→  ", size + 1, col, True); run(p, txt, size, INK, bold)


# ───────────── 슬라이드 ─────────────
# 0. 표지
s = slide()
bar = s.shapes.add_shape(1, 0, Inches(2.35), SW, Inches(2.55))
bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
_, tf = tb(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(2.1), MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]; run(p, "HSMM 레짐 모델 — 2차 보고", 34, WHITE, True)
p2 = tf.add_paragraph(); run(p2, "박사님 피드백 반영 & Risk Overlay(vol-targeting) 검증", 18, RGBColor(0xCF, 0xDA, 0xEC))
_, tf2 = tb(s, Inches(0.9), Inches(5.15), Inches(11), Inches(0.6))
p = tf2.paragraphs[0]
run(p, "핵심: 지적사항 전부 반영 · FCF 전략에서 MDD 유의 개선(P≥0.97) · vol-targeting은 레짐신호와 99% 중복 → 3안 검토", 13, GRAY)

# 1. 피드백 → 반영
s = slide(); header(s, "①", "지난 피드백 → 반영", "회의(2026-07-21) 지적사항 전부 반영 — 모델 목적을 'MDD 하방 방어'로 재정렬")
fb = [("박사님 지적", "반영 내용", "상태"),
      ("모델 목적 = MDD 최소화·리스크 오버레이", "익스포저 조절기로 재정렬, MDD 중심 평가", "반영"),
      ("emission / transition 피처 분리", "상태=breadth·신저가·추세 / 전환=환율·외국인", "반영"),
      ("Ledoit-Wolf 슈링키지", "emission 공분산에 적용(소표본·고차원 안정)", "반영"),
      ("웜스타트 + 5년창 + 시간감쇠", "연1회 재적합·나머지 필터만, HL 36월", "반영"),
      ("Viterbi → filtered 연속확률", "forward 필터 P(bear) → 연속 익스포저", "반영"),
      ("환율: 변화율·Z·레벨가중", "Δ3M × (레벨 / 1년평균), Z-score", "반영"),
      ("금리·EPS 제외", "한국 금리 무효·수출주 환율착시 → 제외", "반영"),
      ("vol-targeting(목표/현재변동성, 하한0.2)", "20일실현 or 60일하방, 하한 0.2 — 본 검증 대상", "반영 ★"),
      ("Z단계 리밸밴드", "0.15 밴드·0.05 스텝(회전·휩소 ↓)", "반영")]
table(s, fb, Inches(0.55), Inches(1.5), Inches(12.25), [5, 5.3, 1.4], header_h=0.4, row_h=0.5, fs=11.5)
_, tf = tb(s, Inches(0.55), Inches(6.5), Inches(12.25), Inches(0.8))
p = tf.paragraphs[0]
run(p, "※ Seed robustness: ", 12, NAVY, True)
run(p, "8개 seed 중 7개가 동일 해로 수렴 → 지배적 해로 고정(재현성 확보). 일부 초기화는 작은 cold-start 창(13개월)發 국소해로 이탈 → "
        "실무 적용 시 multi-start initialization(가중 likelihood 기준 best 선택)으로 근본 대응.", 12, INK)

# 2. 모델 구조
s = slide(); header(s, "②", "반영된 모델 구조", "상태는 시장 내부건전성, 전환은 스트레스가 변조 → filtered 연속확률 → 익스포저")
bullets(s, [("[상태] Emission — '지금 강세/약세인가'", "", NAVY, True),
            ("·", "breadth(200일선 위 종목비율) · 52주 신저가비율 · 추세 log(K/MA200)", GRAY, False),
            ("·", "2-state Gaussian + Ledoit-Wolf · 5년 롤링창 + 시간감쇠", GRAY, False),
            (" ", "", GRAY, False),
            ("[전환] Transition — '상태가 바뀔까'", "", NAVY, True),
            ("·", "환율 Δ3M(레벨가중) + 외국인 3M 순매수 → 스트레스 Z", GRAY, False),
            ("·", "스트레스가 duration exit-hazard를 변조(입력의존 전이)", GRAY, False)],
        Inches(0.55), Inches(1.55), Inches(6.6), Inches(4.5), size=13, gap=4)
box = s.shapes.add_shape(1, Inches(7.4), Inches(1.7), Inches(5.4), Inches(3.9))
box.fill.solid(); box.fill.fore_color.rgb = ALT; box.line.color.rgb = NAVY; box.line.width = Pt(1)
_, tf = tb(s, Inches(7.65), Inches(1.9), Inches(4.95), Inches(3.5))
for i, (mk, txt) in enumerate([("출력", "filtered P(bear) → EMA 스무딩"),
                               ("익스포저", "exposure = clip[ (1−Pbear) · vol_factor , 0.2, 1 ]"),
                               ("vol_factor", "1 − Pbear · cut"),
                               ("cut", "1 − min(1, σ_target / σ_now)"),
                               ("σ_target", "역사적 평균 변동률(인과적 확장평균)"),
                               ("리밸", "Z단계 밴드 0.15 · 0.05 스텝")]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(9)
    run(p, f"{mk}  ", 12.5, BLUE, True); run(p, txt, 12.5, INK)

# 3. 검증 설계
s = slide(); header(s, "③", "검증 설계", "지수가 아닌 '실제 FCF 전략'에 연속 익스포저를 씌워 검증")
bullets(s, ["대상: FCF_YIELD추가전략 (KOSPI 월간), 2018-04 ~ 2026-07, 100개월",
            "벤치마크: FCF불 전략 단독 (익스포저 = 1)",
            "오버레이: overlay_ret[t] = exposure[t] × FCF수익[t], 나머지 현금(0)",
            "인과 타이밍: M월 수익 × M−1월말 확정 익스포저 (lookahead 없음)",
            "비교군: A(20일 실현변동성) · B(60일 하방변동성) · pbear만(vol 無)",
            "목표변동률 = 각 측정치의 인과적 확장평균(=역사적 평균, 박사님 스펙)",
            "유의성: stationary block bootstrap 5000회(블록 6M, paired) → 차이 90%CI · P(우세)"],
        Inches(0.55), Inches(1.6), Inches(12.2), Inches(4.8), size=14.5, gap=10)

# 4. 결과 — MDD
s = slide(); header(s, "④", "결과 ① — MDD 유의미 개선", "통계적으로 견고한 효과 = MDD 방어 (박사님 목적 그대로 달성)")
mt = [("전략", "CAGR", "Sharpe", "MDD", "Calmar", "평균exp"),
      ("FCF불 단독(BM)", "19.2%", "0.74", "−43.7%", "0.44", "1.00"),
      ("A: 20일 실현변동성", "18.1%", "0.84", "−34.5%", "0.52", "0.64"),
      ("B: 60일 하방변동성", "17.9%", "0.84", "−32.5%", "0.55", "0.63"),
      ("pbear만 (vol 無)", "19.5%", "0.84", "−33.7%", "0.58", "0.66")]
table(s, mt, Inches(0.55), Inches(1.55), Inches(6.5), [3, 1.4, 1.4, 1.5, 1.4, 1.4], hl=4, row_h=0.5, fs=11)
s.shapes.add_picture(str(ASSET / "mdd.png"), Inches(7.3), Inches(1.5), height=Inches(3.0))
conclusion(s, [("MDD: 오버레이 3종 모두 BM 대비 +12~13%p, P=0.97~0.99 → 90%CI가 0 제외 = 유의", GREEN, True),
               ("Sharpe·CAGR·Calmar 개선은 90%CI가 0 포함 → 100개월 표본으론 미입증(부수 효과)", GRAY, False),
               ("MDD 병목 이동: BM 2018→2020-03(코로나) / 오버레이 2021-05→2022-12(2022 하락장)", BLUE, False)],
           Inches(4.95))

# 5. vol-targeting 유무
s = slide(); header(s, "⑤", "결과 ② — vol-targeting 유무는 무의미", "분모(20일실현 vs 60일하방)도, vol-targeting 유무도 방어 성과는 통계적으로 동일")
vt = [("비교 (paired bootstrap)", "핵심 차이 Δ [90%CI]", "P(우세)", "판정"),
      ("오버레이 vs BM", "MDD +12~13%p [+1.0, +29.1]", "0.97~0.99", "유의 (방어 확실)"),
      ("A vs B", "MDD −0.9%p [−2.4, +0.4]", "0.13~0.57", "무의미 (동일)"),
      ("pbear vs A", "CAGR +1.2%p [+0.0, +3.7]", "0.96", "pbear 소폭 우세(CAGR만)"),
      ("pbear vs A·B", "MDD·Sharpe·Calmar", "≈ 0.5", "동급 (방어 차이 無)")]
table(s, vt, Inches(0.55), Inches(1.6), Inches(12.25), [3.3, 4.2, 1.6, 3], hl=1, row_h=0.55, fs=12)
conclusion(s, [("분모 종류(20일실현/60일하방) 선택은 통계적으로 무의미 — 결론 안 바뀜", INK, True),
               ("vol-targeting을 얹어도 방어(MDD·Sharpe) 개선 0, CAGR만 소폭(−1.2%p) 손해", RED, True)],
           Inches(4.6))

# 6. 왜 추가 효과가 제한적이었나
s = slide(); header(s, "⑥", "왜 vol-targeting의 추가 효과가 제한적이었는가", "vol-targeting이 레짐 신호와 대부분 겹쳐, 새로 더해줄 정보가 적었다")
bullets(s, [("① 레짐 신호가 이미 위험을 반영", "시장이 위험해지면 pbear가 오르고 주식 익스포저가 자동으로 줄어듦", NAVY, True),
            ("② 작동 구간이 겹침", "변동성이 커지는 때 = pbear가 이미 높은 때 → vol-targeting은 같은 위험을 한 번 더 줄이는 셈", NAVY, True),
            ("③ 결과: 거의 같은 포트폴리오", f"두 방식의 익스포저 경로가 상관 {corr:.2f}로 거의 동일", NAVY, True),
            (" ", "→ vol-targeting은 미세하게 더 깎아 CAGR만 소폭 손해, MDD 방어는 그대로", GRAY, False)],
        Inches(0.55), Inches(1.75), Inches(6.4), Inches(3.9), size=13.5, gap=14)
s.shapes.add_picture(str(ASSET / "exp.png"), Inches(7.15), Inches(1.6), height=Inches(3.15))
_, tf = tb(s, Inches(0.55), Inches(6.15), Inches(12.2), Inches(1.0))
p = tf.paragraphs[0]
run(p, "※ 단, vol을 레짐 신호에 곱한 현재 구현 기준. 레짐과 독립적인 vol-targeting은 별도 검증 필요.", 12, RED, italic=True)

# 7. 질문 (3안)
s = slide(); header(s, "⑦", "질문 — vol-targeting을 어떻게 가져갈까", "MDD 방어는 확정. vol-targeting 채택 여부에 대한 3가지 선택지")
op = [("A안 · 채택 보류", "pbear-only. 같은 MDD 방어 + CAGR 소폭↑ · 가장 단순(오컴). 목표변동률·분모 파라미터 제거", GREEN),
      ("B안 · 조건부 유보", "현 '소프트 비대칭' 구현은 pbear와 99% 중복. 순수 vol-targeting부터 검증 후 결정", BLUE),
      ("C안 · 스펙 유지", "유의 개선 없어도 실전 정합·표본 밖 강건성 위해 20일실현 vol-targeting(목표 18%) 유지", NAVY)]
_, tf = tb(s, Inches(0.55), Inches(1.5), Inches(12.25), Inches(2.7))
for i, (t, d, c) in enumerate(op):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(10)
    run(p, f"{t}   ", 15, c, True); run(p, d, 13, INK)
ln = s.shapes.add_shape(1, Inches(0.55), Inches(4.35), Inches(12.25), Pt(1.2))
ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC); ln.line.fill.background()
bullets(s, [("Q1 ★", "세 방향 중 어느 쪽? — 하방방어가 목적이면 A안(단순)이 타당한가?", RED, True),
            ("Q2", "순수 vol-targeting(pbear 독립 스케일링)을 별도로 검증할 가치가 있나?", NAVY, True),
            ("Q3", "목표변동률을 실전처럼 고정(18%)으로 vs 인과적 확장평균 — 어느 쪽이 정석인가?", NAVY, True),
            ("Q4", "표본 확장(2005~, 지수기반 피처)으로 Sharpe 유의성까지 재검할 가치가 있나?", NAVY, True)],
        Inches(0.55), Inches(4.55), Inches(12.2), Inches(2.6), size=13.5, gap=8)

prs.save(str(OUT))
print(f"저장 완료: {OUT}")
print(f"익스포저 corr(A,pbear) = {corr:.3f}")
