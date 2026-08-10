# -*- coding: utf-8 -*-
"""레짐 모델 탐색과정 PPT (3장) — 쉬운 말 버전, 한글 차트."""
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm
import numpy as np
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# 한글 폰트
for cand in ["AppleGothic", "Apple SD Gothic Neo", "NanumGothic", "Malgun Gothic"]:
    if any(cand in f.name for f in fm.fontManager.ttflist) or cand == "AppleGothic":
        plt.rcParams["font.family"] = cand; break
plt.rcParams["axes.unicode_minus"] = False
CH = Path("analysis/ppt_charts"); CH.mkdir(exist_ok=True)
GREEN = "#1B6E3C"; RED = "#B12A2A"; BLUE = "#2E5B9A"; GRAY = "#888888"; ORANGE = "#C77A14"

# C1: 랜덤 초기값별 약세장 구분 점수
seeds = ["①", "②", "③", "④", "⑤"]
no_fx = [-0.93, -1.75, -2.18, 1.75, -1.91]; fx = [1.84, 1.53, 1.53, 1.36, 1.53]
x = np.arange(5); w = .38
fig, ax = plt.subplots(figsize=(6.6, 4.0))
ax.bar(x-w/2, no_fx, w, label="내부지표만", color=RED, alpha=.85)
ax.bar(x+w/2, fx, w, label="+ 환율", color=GREEN, alpha=.9)
ax.axhline(0, color="k", lw=.8)
ax.set_xticks(x); ax.set_xticklabels(seeds); ax.set_xlabel("컴퓨터 랜덤 초기값 (5가지로 돌려봄)")
ax.set_ylabel("약세장 구분 점수 (높을수록 잘 구분)")
ax.set_title("환율이 '운빨'을 없애줌\n내부지표만: 초기값 따라 +/- 뒤집힘 → 환율 넣으니 항상 +", fontsize=10.5)
ax.legend(fontsize=10); fig.tight_layout(); fig.savefig(CH/"k1.png", dpi=150); plt.close()

# C2: 추가 정보별 성적
cfgs = ["기본\n(환율)", "+금리", "+달러", "+뉴스", "3단계", "4단계"]
mu = [1.56, 2.32, 1.77, -0.80, -0.58, 0.43]; sd = [0.16, 0.12, 0.0, 3.29, 1.82, 1.90]
note = ["채택", "수익엔\n도움X", "약함", "모델\n망가짐", "불안정", "불안정"]
cols = [GREEN, ORANGE, BLUE, RED, RED, RED]
fig, ax = plt.subplots(figsize=(6.6, 4.0))
ax.bar(np.arange(6), mu, yerr=sd, color=cols, alpha=.85, capsize=4)
ax.axhline(0, color="k", lw=.8)
for i,(m,s,t) in enumerate(zip(mu,sd,note)):
    ax.text(i, m+(s+0.25 if m>=0 else -s-0.6), t, ha="center", fontsize=8.5, color=cols[i])
ax.set_xticks(np.arange(6)); ax.set_xticklabels(cfgs, fontsize=9)
ax.set_ylabel("약세장 구분 점수 (막대=평균, 선=흔들림)")
ax.set_title("정보를 더 넣어봤지만…\n기본(환율)만 안정적으로 좋고, 나머진 흔들리거나 망가짐", fontsize=10.5)
fig.tight_layout(); fig.savefig(CH/"k2.png", dpi=150); plt.close()

# C3: vs ai_v2
fig, axs = plt.subplots(1, 2, figsize=(7.2, 3.8)); xx = np.arange(2); w = .36
axs[0].bar(xx-w/2, [7,1.53], w, label="우리모델", color=GREEN); axs[0].bar(xx+w/2, [6,0.23], w, label="기존 AI", color=GRAY)
axs[0].set_xticks(xx); axs[0].set_xticklabels(["위기 잡은 수\n(7번중)","약세장\n구분점수"], fontsize=9)
axs[0].set_title("위기 잘 잡나? → 우리가 우위", fontsize=10); axs[0].legend(fontsize=9)
axs[1].bar(xx-w/2, [353,38.1], w, label="우리모델", color=GREEN); axs[1].bar(xx+w/2, [406,40.6], w, label="기존 AI", color=GRAY)
axs[1].set_xticks(xx); axs[1].set_xticklabels(["누적수익\n(%)","최대낙폭\n(%)"], fontsize=9)
axs[1].set_title("돈은? → AI가 더 범", fontsize=10); axs[1].legend(fontsize=9)
fig.tight_layout(); fig.savefig(CH/"k3.png", dpi=150); plt.close()

# ===== PPT =====
NAVY=RGBColor(0x1F,0x33,0x55); BLUEc=RGBColor(0x2E,0x5B,0x9A); GR=RGBColor(0x44,0x44,0x44)
REDc=RGBColor(0xB1,0x2A,0x2A); GRNc=RGBColor(0x1B,0x6E,0x3C)
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)

def slide(title):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    tf=s.shapes.add_textbox(Inches(0.45),Inches(0.3),Inches(12.4),Inches(0.9)).text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=title; p.font.size=Pt(26); p.font.bold=True; p.font.color.rgb=NAVY
    return s
def body(s, items, left, top, width, size=14):
    tf=s.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(5.8)).text_frame; tf.word_wrap=True
    for i,(lvl,txt,col,bold) in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=txt; p.level=lvl; p.font.size=Pt(size-lvl*1.5); p.font.bold=bold; p.font.color.rgb=col; p.space_after=Pt(5)

# Slide 1
s=slide("① 정보를 더 넣어봤지만, 다 실패했습니다")
body(s,[
 (0,"우리 모델에 '약세장을 더 잘 맞히려고' 여러 정보를 하나씩 넣어봤습니다.",GR,False),
 (0,"금리 / 달러 / 뉴스 / 구리·금 …",NAVY,True),
 (0,"결과:",NAVY,True),
 (1,"금리·달러 → 약세장 구분은 좋아졌는데, 정작 '돈'은 더 못 벌었습니다",REDc,False),
 (1,"뉴스 → 넣었더니 모델이 들쭉날쭉 망가졌습니다 (오른쪽 빨강)",REDc,False),
 (1,"구리·금 → 새 정보긴 한데 약세장 예측엔 도움 안 됐습니다",REDc,False),
 (0,"→ 한 줄 결론: 정보를 더 넣는다고 좋아지지 않았다.",GRNc,True),
 (1,"\"약세장 구분을 잘한다\"와 \"실제로 돈을 번다\"는 다른 얘기였습니다.",GR,False),
], 0.5, 1.5, 6.1)
s.shapes.add_picture(str(CH/"k2.png"), Inches(6.9), Inches(1.9), width=Inches(6.1))

# Slide 2
s=slide("② '환율' 하나가 모델을 운빨에서 구해줬습니다")
body(s,[
 (0,"문제: 시장 내부지표(오르는 종목이 얼마나 많나)만 쓰면,",GR,False),
 (1,"컴퓨터 랜덤 초기값에 따라 결과가 뒤집혔습니다 = '운빨' (오른쪽 빨강 막대)",REDc,False),
 (0,"환율을 넣었더니:",NAVY,True),
 (1,"어떤 초기값으로 돌려도 항상 약세장을 잘 구분 (오른쪽 초록)",GRNc,True),
 (1,"즉 환율은 '예측을 잘해서'가 아니라 '결과를 안 흔들리게' 해줘서 채택",GR,False),
 (0,"이것저것 더 바꿔봤지만:",NAVY,True),
 (1,"단계를 3·4개로 쪼개기 → 더 불안정해짐",REDc,False),
 (1,"예측 방식 바꾸기 → 더 빨리 잡지도 못함",REDc,False),
 (0,"→ 가장 단순한 '2단계 + 환율'이 최선이었습니다.",GRNc,True),
], 0.5, 1.5, 6.0)
s.shapes.add_picture(str(CH/"k1.png"), Inches(6.9), Inches(2.0), width=Inches(6.1))

# Slide 3
s=slide("③ 위기는 우리가 더 잘 잡는데, 돈은 AI가 더 벌었습니다")
body(s,[
 (0,"우리 모델 vs 기존 AI(뉴스 읽는 모델):",NAVY,True),
 (1,"위기 감지: 우리가 우위 (7번 중 7번 vs 6번), 낙폭도 더 잘 막음",GRNc,False),
 (1,"하지만 최종 수익은 AI가 더 높았습니다 (406% vs 353%)",REDc,False),
 (0,"왜 AI가 돈을 더 벌었나? 뜯어보니:",NAVY,True),
 (1,"2022 하락장: AI는 뉴스 보고 계속 '약세' 유지 / 우리는 잠깐 반등에 '강세'로 깜빡여 손해",GR,False),
 (1,"반등장: 우리는 '약세'로 보고 손절했다가 바로 튀어오르는 걸 놓침 / AI는 버텨서 먹음",GR,False),
 (0,"→ AI가 이긴 건 '더 좋은 모델'이라서가 아니라 '뉴스'라는 정보가 있어서.",NAVY,True),
 (1,"그 뉴스 정보를 우리 모델에 넣어보려 했지만 모두 실패했습니다.",REDc,False),
 (0,"다음 계획:",NAVY,True),
 (1,"약세장에 전략을 통째로 바꾸지 말고, 그냥 '비중만 줄이기' → 손절 손해부터 막기",BLUEc,True),
], 0.5, 1.45, 6.4)
s.shapes.add_picture(str(CH/"k3.png"), Inches(7.1), Inches(2.5), width=Inches(5.9))

out=Path("레짐모델_쉽게정리.pptx"); prs.save(out)
print(f"저장: {out.resolve()} (3장)")
