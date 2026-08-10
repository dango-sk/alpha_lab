"""레짐 예측 모델 연구과정 정리 PPT (3장) 생성 — 박사님 보고용."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pathlib import Path

NAVY = RGBColor(0x1F, 0x33, 0x55); BLUE = RGBColor(0x2E, 0x5B, 0x9A)
GRAY = RGBColor(0x55, 0x55, 0x55); RED = RGBColor(0xB1, 0x2A, 0x2A); GREEN = RGBColor(0x1B, 0x6E, 0x3C)
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)


def slide(title, sub):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0)).text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]; p.text = title; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = NAVY
    p2 = tb.add_paragraph(); p2.text = sub; p2.font.size = Pt(13); p2.font.color.rgb = BLUE; p2.font.italic = True
    return s


def body(s, items, top=1.5, left=0.6, width=12.1, size=13):
    tf = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.6)).text_frame
    tf.word_wrap = True
    for i, (lvl, txt, col, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.level = lvl
        p.font.size = Pt(size - lvl * 1); p.font.bold = bold
        p.font.color.rgb = col; p.space_after = Pt(3)


# ===== Slide 1 =====
s = slide("KOSPI 레짐 예측 모델 ① — 문제의식과 핵심 고민",
          "월단위 강세/약세 레짐 → FCF 팩터전환 전략(강세전략↔약세전략) 스위칭")
body(s, [
 (0, "■ 출발 문제", NAVY, True),
 (1, "기존 LLM 멀티모델(ai_v2)이 약세장 Recall ≈ 0% — 위기를 거의 못 잡음", GRAY, False),
 (0, "■ 고민 1 — \"무엇을 예측할 것인가\"", NAVY, True),
 (1, "다음 달 방향(↑/↓)은 13개 모델군 모두 AUC ≈ 0.5 → 예측 불가 확인", GRAY, False),
 (1, "→ 목표를 '방향'이 아닌 '위험 국면(전환점) 탐지'로 재정의", GREEN, True),
 (1, "→ 평가도 방향 정확도 폐기, Event Recall·Lead·False Alarm·Whipsaw·Bull/Bear 수익격차로", GRAY, False),
 (0, "■ 고민 2 — \"어떤 모델 구조\"", NAVY, True),
 (1, "하드코딩 룰 거부(과적합·자의성) → 비지도 상태추론 HMM 선택", GRAY, False),
 (1, "Full-Cov HMM은 분별력 좋으나 전환 잦음(Whipsaw↑) → 안정화 4종 비교(Sticky/Penalty/Hysteresis/HSMM)", GRAY, False),
 (1, "→ HSMM(상태 지속기간을 모델이 학습) 채택: 분별력 유지하며 전환 안정", GREEN, True),
 (0, "■ 고민 3 (결정적) — \"왜 결과가 seed(초기값)마다 다른가\"", NAVY, True),
 (1, "시장 내부지표(breadth 등)만으론 Full-Cov HMM이 초기값에 취약 — 분별력이 seed 운에 좌우", RED, False),
 (1, "→ 환율(원화 3개월 변화) 추가 → 5개 seed 전부 격차 +1.5%p로 안정화", GREEN, True),
 (1, "환율의 역할 = '예측력'이 아니라 'robustness' : 시장 내부에 자본흐름(외부) 독립 축을 더해 상태 식별을 고정", GRAY, False),
])

# ===== Slide 2 =====
s = slide("KOSPI 레짐 예측 모델 ② — 탐색 과정과 모델 선택 이유",
          "변수도 추가하고, HMM 구조·예측방식도 바꿔봤지만 — 모두 검증 후 기각")
body(s, [
 (0, "■ 추가·변경 시도 → 결과 → 채택 여부", NAVY, True),
 (1, "[피처] 금리(US10Y)·DXY : 분별력 +1.53→+2.41↑ 이나 FCF 수익 무익 → 기각", GRAY, False),
 (1, "[피처] LLM 뉴스(Tightening/Stress) : 파일럿 설명력 있으나 raw 피처 투입 시 robustness 파괴·2023 과잉방어 → 기각", GRAY, False),
 (1, "[피처] 구리/금 비율 : 통계적 독립(77%)이나 bear 예측 증분가치 0·breadth보다 후행 → 기각", GRAY, False),
 (1, "[구조] 3/4-state(Bull/SlowBear/FastBear/Recovery) : 상태는 의미있게 분화되나 Bull/Bear 매핑 시 격차 음수·불안정 → 기각", GRAY, False),
 (1, "[예측방식] 전이행렬 1-step : Lead 개선 없음(+0.7 vs +0.9), FA·Whip 소폭↑ → 기각", GRAY, False),
 (0, "■ 반복 확인된 교훈", NAVY, True),
 (1, "분별력(Bull/Bear 수익격차) ≠ FCF 실제 수익  —  분별력을 올려도 수익으로 안 이어짐", RED, True),
 (1, "독립 정보 ≠ 유용한 정보  /  피처·구조 추가만으로는 base를 못 이김", RED, False),
 (0, "■ 그래서 최종 선택", NAVY, True),
 (1, "Full-Cov HSMM (2-state) + 환율  —  가장 단순하면서 robust·분별력 최선", GREEN, True),
 (1, "피처 6개: breadth · Δbreadth · 52주 신저가비율 · Δ신저가 · 모멘텀(1m−6m) · 환율Δ3m  (시장내부 5 + 자본흐름 1)", GRAY, False),
 (1, "산출: 5-seed 다수결 consensus · walk-forward · 미래정보 누수 없음", GRAY, False),
])

# ===== Slide 3 =====
s = slide("KOSPI 레짐 예측 모델 ③ — 결과 · 한계 · 다음 단계",
          "vs 기존 챔피언 ai_v2 (overlap 2018~2026)")
body(s, [
 (0, "■ 성능 (hsmm+환율  vs  ai_v2)", NAVY, True),
 (1, "Event Recall 7/7 vs 6/7  |  Bull/Bear 격차 +1.53%p vs +0.23%p  |  MDD −38% vs −41%   → 위기탐지·분별력·낙폭방어 우위", GREEN, False),
 (1, "FCF 누적수익 329~353% vs 406%   → 절대수익은 ai_v2 우위(아직 못 이김)", RED, False),
 (0, "■ ai_v2가 수익이 더 높았던 이유 (월별 attribution)", NAVY, True),
 (1, "① 2022 긴축 베어: ai는 뉴스/매크로로 계속 Bear 유지, hsmm은 breadth가 반등마다 Bull로 깜빡여 하락 구간을 Bull로 맞음 (−7.4%p)", GRAY, False),
 (1, "② 회복장: hsmm은 Bear로 빠져 타이트한 손절이 바닥에서 매도(반등 놓침), ai는 Bull 유지로 흡수 (−5.2%p, 손절 139건)", GRAY, False),
 (1, "③ ai는 Bear월·Whipsaw 적어 거래·손절 드래그 작음   → 종목선택·분별력 아니라 '정보·타이밍' 차이", GRAY, False),
 (0, "■ 핵심 한계", NAVY, True),
 (1, "ai_v2의 우위는 본질적으로 '뉴스/매크로 정보' — 이를 HMM 피처로 주입하는 시도는 모두 실패(분별력만↑ or 모델 파괴)", RED, True),
 (1, "즉 레짐 모델 단독 개선으로는 ai_v2 수익 갭(≈53%p)을 닫기 어려움", RED, False),
 (0, "■ 다음 단계", NAVY, True),
 (1, "Risk Overlay: 약세전략 전면교체 대신 강세전략(alpha engine) 유지 + Bear월 익스포저만 축소 → 회복장 손절드래그 정조준", BLUE, True),
 (1, "+ ai_v2 하이브리드 게이트: 매크로 지속성으로 hsmm의 2022 깜빡임 보완 (정보를 피처가 아니라 레짐 결합으로)", BLUE, False),
])

out = Path("레짐예측모델_연구정리.pptx")
prs.save(out)
print(f"저장 완료: {out.resolve()}  ({len(prs.slides.__iter__.__self__._sldIdLst)}장)")
