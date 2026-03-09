"""
퀀트 모델 리뷰 PPT 생성
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── 색상 팔레트 (미니멀 화이트) ───
BG_DARK = RGBColor(0xFF, 0xFF, 0xFF)       # 흰 배경
BG_CARD = RGBColor(0xF5, 0xF5, 0xF5)       # 카드: 연한 회색
BG_ACCENT = RGBColor(0xEB, 0xEB, 0xEB)     # 테이블 짝수행
BLUE = RGBColor(0x1A, 0x56, 0xDB)          # 포인트 블루
GREEN = RGBColor(0x05, 0x7A, 0x55)         # 긍정 (짙은 초록)
RED = RGBColor(0xC8, 0x1E, 0x1E)           # 부정 (짙은 빨강)
YELLOW = RGBColor(0xB4, 0x5D, 0x09)        # 강조 (짙은 노랑/갈색)
WHITE = RGBColor(0x1F, 0x1F, 0x1F)         # 본문 텍스트 (거의 검정)
GRAY = RGBColor(0x6B, 0x72, 0x80)          # 보조 텍스트
LIGHT_GRAY = RGBColor(0x9C, 0xA3, 0xAF)    # 약한 보조 텍스트


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=WHITE, bullet_color=BLUE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "맑은 고딕"
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def add_table(slide, left, top, width, rows_data, col_widths=None, font_size=11):
    """rows_data: list of lists. First row = header."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    row_height = 0.32 if n_rows > 10 else 0.4
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                          Inches(width), Inches(row_height * n_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "맑은 고딕"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = WHITE
                paragraph.alignment = PP_ALIGN.CENTER

            # 배경색
            fill = cell.fill
            fill.solid()
            if r == 0:
                fill.fore_color.rgb = BLUE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            elif r % 2 == 0:
                fill.fore_color.rgb = BG_CARD
            else:
                fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    return table_shape


def add_accent_box(slide, left, top, width, height, text, font_size=12, bg=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(0.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(0x1F, 0x1F, 0x1F)
    p.font.name = "맑은 고딕"
    return shape


def slide_title(slide, title, subtitle=None):
    set_slide_bg(slide)
    # 상단 파란 라인
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.5), Inches(0.4), Inches(1.5), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    add_text_box(slide, 0.5, 0.55, 9, 0.6, title, font_size=24, bold=True, color=RGBColor(0x1F, 0x1F, 0x1F))
    if subtitle:
        add_text_box(slide, 0.5, 1.05, 9, 0.4, subtitle, font_size=14, color=GRAY)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank

    # ════════════════════════════════════════
    # 표지
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, 1, 2.0, 11, 1.0, "퀀트 모델 리뷰 & 회귀 팩터 분석",
                 font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 3.2, 11, 0.6, "이현자산운용 실무진 미팅 자료",
                 font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.2, 11, 0.5, "2026년 3월",
                 font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    # 하단 라인
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(4), Inches(3.9), Inches(5.33), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    # ════════════════════════════════════════
    # 슬라이드 1: 목적
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "모델의 목적", "왜 만들었고, 무엇을 하는 모델인가")

    add_text_box(slide, 0.5, 1.6, 5, 0.4, "왜 만들었나", font_size=18, bold=True, color=BLUE)
    add_bullet_list(slide, 0.7, 2.1, 5.5, 2.0, [
        "• 이현의 밸류 투자 프로세스를 정량적 모델로 체계화",
        "• 종목 선정 시 주관적 판단의 편향을 줄이고, 일관된 기준 적용",
        "• 백테스트를 통해 전략의 유효성을 사전 검증",
    ], font_size=14, color=WHITE)

    add_text_box(slide, 7, 1.6, 5, 0.4, "무엇을 하는 모델인가", font_size=18, bold=True, color=BLUE)
    add_bullet_list(slide, 7.2, 2.1, 5.5, 2.0, [
        "• 매월 1회, KOSPI 유니버스에서 상위 30개 종목 자동 선정",
        "• 선정 기준: 밸류 + 회귀 + 성장 + 차별화",
        "• 종목별 비중은 시총 비례 + 최대 10% 캡",
    ], font_size=14, color=WHITE)

    add_accent_box(slide, 0.5, 4.2, 12.3, 1.8,
                   "핵심: '어떤 종목이 저평가인가'를 4가지 관점에서 점수화하여\n"
                   "매월 최적의 30개 종목 포트폴리오를 구성하는 시스템",
                   font_size=15)

    # ════════════════════════════════════════
    # 슬라이드 2: 데이터 소스
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "데이터 소스 & 유니버스", "어디서 데이터를 가져오고, 어떤 종목을 대상으로 하는가")

    add_text_box(slide, 0.5, 1.6, 5, 0.4, "데이터 소스", font_size=18, bold=True, color=BLUE)
    add_table(slide, 0.5, 2.1, 5.5, [
        ["데이터", "소스", "내용"],
        ["재무제표", "FnSpace", "PER, PBR, ROE, ROIC 등"],
        ["컨센서스", "FnSpace", "Forward EPS, Forward PER"],
        ["주가/시총", "FnSpace", "일별 종가, 시가총액, 거래량"],
        ["벤치마크", "KOSPI 200 ETF", "비교 기준 수익률"],
    ], font_size=12)

    add_text_box(slide, 7, 1.6, 5, 0.4, "유니버스 정의", font_size=18, bold=True, color=BLUE)
    add_table(slide, 7, 2.1, 5.8, [
        ["항목", "설정"],
        ["시장", "KOSPI만 (KOSDAQ 제외)"],
        ["시총 하한", "5천억 원"],
        ["유니버스 시총 범위", "약 5천억 ~ 500조 (삼성전자 등)"],
        ["유니버스 규모", "약 680개 → 상위 30개 선정"],
        ["제외 업종", "금융업, SPAC, ETF, REIT"],
        ["퀄리티 필터", "부채비율 200% 이하, 적자 1Q 이하, 거래대금 5억+"],
    ], font_size=12)

    # ════════════════════════════════════════
    # 슬라이드 3: 모델 구조
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "모델 구조 — 4가지 팩터", "총점 = 밸류(35%) + 회귀(30%) + 성장(20%) + 차별화(15%)")

    add_table(slide, 0.5, 1.8, 12.3, [
        ["팩터", "비중", "보는 것", "지표"],
        ["밸류", "35%", "현재 주가가 싼가", "PER, PBR, EV/EBITDA, PCF (현재+Forward)"],
        ["회귀", "30%", "적정가 대비 괴리가 있나", "PBR-ROE, EV/IC-ROIC, F.PER-이익성장, F.EV/EBIT-EBIT성장 회귀"],
        ["성장", "20%", "실적이 성장하고 있나", "매출 성장률 (현재+Forward)"],
        ["차별화", "15%", "시장 기대가 올라가고 있나", "EPS 컨센서스 변화율 (모멘텀)"],
    ], col_widths=[1.2, 0.8, 3.0, 7.3], font_size=13)

    add_accent_box(slide, 0.5, 4.5, 12.3, 2.2,
                   "회귀 팩터란?\n\n"
                   "단순히 '싸다'가 아니라 '같은 수준의 수익성(ROE) 대비 싸다'를 측정\n\n"
                   "예: ROE 15%인 A종목(PBR 1.0)과 B종목(PBR 0.7)이 있을 때\n"
                   "→ B가 회귀선 아래 → '적정가보다 저평가' → 매력적",
                   font_size=14)

    # ════════════════════════════════════════
    # 슬라이드 4: 백테스트 설정
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "백테스트 설정")

    add_table(slide, 0.5, 1.8, 6, [
        ["항목", "설정"],
        ["기간", "2021.01 ~ 2026.03 (약 5년)"],
        ["리밸런싱", "월 1회"],
        ["종목 수", "30개"],
        ["비중 방식", "시총 비례 + 최대 10% 캡"],
        ["거래비용", "편도 30bp"],
        ["벤치마크", "KOSPI 200"],
    ], col_widths=[2.0, 4.0], font_size=14)

    # 비중 산출 방식 상세
    add_text_box(slide, 7, 1.6, 5.8, 0.4, "비중 산출 방식", font_size=18, bold=True, color=BLUE)
    add_accent_box(slide, 7, 2.1, 5.8, 3.5,
                   "① 초기 비중 = 각 종목 시총 / 30종목 시총 합계\n"
                   "   → 시총이 큰 종목일수록 비중이 높음\n\n"
                   "② 개별 종목 비중 상한 10% 적용\n"
                   "   → 10% 초과 종목의 초과분을 나머지에 재배분\n"
                   "   → 수렴할 때까지 반복 (최대 20회)\n\n"
                   "③ 결과 예시 (30종목 기준)\n"
                   "   → 균등이면 3.3%씩 / 실제는 1~10% 분포\n"
                   "   → 삼성전자 등 대형주가 10%에 캡핑됨\n\n"
                   "④ 회전율 = (편입+편출) / (2×종목수)\n"
                   "   → 매월 리밸런싱 시 교체된 종목 비율",
                   font_size=12, bg=BG_CARD)

    # ════════════════════════════════════════
    # 슬라이드 5: 단일 팩터 롱숏 검증
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "단일 팩터 롱숏 검증", "각 팩터가 단독으로 종목 선별력이 있는가 (5분위 Q1 롱 - Q5 숏)")

    add_table(slide, 0.3, 1.7, 12.7, [
        ["팩터", "Sharpe", "연환산", "누적 L/S", "t-stat", "p-value", "승률", "판정"],
        ["Forward EPS 모멘텀", "+0.79", "+13.0%", "+73.7%", "1.78", "0.081", "60%", "강력 (준유의)"],
        ["Trailing 매출 성장률", "+0.78", "+12.2%", "+73.0%", "1.79", "0.079", "60%", "강력 (준유의)"],
        ["Forward 매출 성장률", "+0.73", "+13.5%", "+80.1%", "1.68", "0.098", "62%", "강력 (준유의)"],
        ["Forward EBIT Growth", "+0.60", "+8.2%", "+44.6%", "1.37", "0.176", "51%", "양호"],
        ["가격 모멘텀 (3M)", "+0.45", "+10.3%", "+48.3%", "1.04", "0.304", "46%", "양호"],
        ["ATT: EV/IC~ROIC", "+0.43", "+8.9%", "+41.3%", "0.98", "0.332", "56%", "양호"],
        ["ATT: PBR~ROE", "+0.27", "+6.4%", "+21.1%", "0.62", "0.540", "48%", "약함"],
        ["Forward PER", "+0.22", "+4.5%", "+14.5%", "0.51", "0.609", "46%", "약함"],
        ["ATT: F.EV/EBIT~EBITG", "+0.20", "+3.1%", "+10.3%", "0.45", "0.654", "49%", "약함"],
        ["ATT: F.PER~EPSG", "-0.09", "-1.6%", "-15.4%", "-0.19", "0.847", "44%", "역효과"],
        ["PBR", "-0.10", "-2.4%", "-24.0%", "-0.24", "0.813", "41%", "역효과"],
        ["Trailing EV/EBITDA", "-0.09", "-1.9%", "-18.5%", "-0.21", "0.834", "48%", "역효과"],
    ], col_widths=[2.8, 1.0, 1.1, 1.2, 1.0, 1.1, 0.8, 1.8], font_size=10)

    add_accent_box(slide, 0.3, 6.3, 12.7, 0.9,
                   "핵심: EPS 모멘텀·매출 성장률이 가장 강력 (Sharpe 0.7~0.8) | 회귀 4종 중 EV/IC~ROIC만 양호, "
                   "F.PER~EPSG는 역효과 | PBR·EV/EBITDA도 단독으로는 역효과",
                   font_size=11)

    # ════════════════════════════════════════
    # 슬라이드 6: 비중 vs 기여도 갭
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "팩터 비중 vs 실제 기여도 갭", "현재 비중 배분이 최적인가?")

    add_table(slide, 0.3, 1.7, 12.7, [
        ["팩터 그룹", "팩터", "현재 비중", "롱숏 Sharpe", "판정"],
        ["차별화 (15%)", "Forward EPS 모멘텀", "15%", "+0.79", "강력 — 비중 대비 기여도 높음 ↑"],
        ["성장 (20%)", "매출 성장률 (T+F)", "20%", "+0.78", "강력 — 비중 대비 기여도 높음 ↑"],
        ["회귀 (30%)", "ATT: EV/IC~ROIC", "5%", "+0.43", "양호 — 4종 중 유일하게 양호"],
        ["회귀 (30%)", "ATT: PBR~ROE", "5%", "+0.27", "약함"],
        ["회귀 (30%)", "ATT: F.PER~EPSG", "10%", "-0.09", "역효과 ⚠️"],
        ["회귀 (30%)", "ATT: F.EV/EBIT~EBITG", "10%", "+0.20", "약함"],
        ["밸류 (35%)", "Forward PER", "5%", "+0.22", "약함"],
        ["밸류 (35%)", "PBR (T+F)", "10%", "-0.10", "역효과 ⚠️"],
        ["밸류 (35%)", "EV/EBITDA (T+F)", "10%", "-0.09", "역효과 ⚠️"],
    ], col_widths=[1.8, 2.8, 1.2, 1.3, 5.6], font_size=11)

    add_accent_box(slide, 0.3, 5.5, 12.7, 1.5,
                   "핵심 질문: 롱숏에서 강력한 팩터 비중을 올리면 전략이 좋아질까?\n\n"
                   "• EPS 모멘텀(Sharpe 0.79)은 비중 15%에 불과 → 비중 확대?\n"
                   "• PBR/EV/EBITDA는 단독 역효과인데 밸류 35% 중 20%를 차지 → 제거?\n"
                   "• 회귀 30% 중 양호한 건 EV/IC~ROIC(5%)뿐 → 약한 회귀 축소?",
                   font_size=13)

    # ════════════════════════════════════════
    # 슬라이드 7: 비중 조정 시뮬레이션
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "비중 조정 시뮬레이션", "롱숏 결과를 근거로 비중을 바꿔보면?")

    add_table(slide, 0.5, 1.8, 12.3, [
        ["전략", "변경 내용", "누적", "CAGR", "Sharpe", "MDD"],
        ["A0 (현재)", "밸류35+회귀30+성장20+차별화15", "+140.1%", "+18.2%", "0.87", "-34.9%"],
        ["역효과 제거", "PBR·EV/EBITDA·F.PER~EPSG 제거, 비례 재배분", "+121.5%", "+16.4%", "0.74", "-39.3%"],
        ["모멘텀/성장 강화", "EPS모멘텀 25%, 매출성장 25%로 확대", "+101.8%", "+14.3%", "0.70", "-39.6%"],
        ["최적 제안", "역효과+약한회귀 제거 + 모멘텀/성장 확대", "+78.0%", "+11.6%", "0.58", "-48.1%"],
        ["KOSPI 200", "벤치마크", "+107.3%", "+14.9%", "0.70", "-33.9%"],
    ], col_widths=[2.0, 4.0, 1.2, 1.1, 1.1, 1.1], font_size=12)

    add_accent_box(slide, 0.5, 5.0, 5.8, 2.0,
                   "핵심 발견: 역효과 팩터를 빼면 오히려 악화\n\n"
                   "• 역효과 제거 시 Sharpe 0.87→0.74 하락\n"
                   "  (위험 대비 수익 효율이 실질적으로 저하)\n"
                   "• 단일 팩터 롱숏의 역효과 ≠ 멀티팩터에서 불필요\n"
                   "• 14팩터 합산 시 종목 간 변별력을 높이는\n"
                   "  보조 역할 → 제거하면 순위 정밀도 하락",
                   font_size=12, bg=RGBColor(0xEC, 0xFD, 0xF5))

    add_accent_box(slide, 6.8, 5.0, 5.8, 2.0,
                   "결론: 현재 14팩터 조합이 최적\n\n"
                   "• 비중 조정으로는 개선 불가 → 구조 자체가 최적\n"
                   "• 개선 방향은 팩터의 '질적 보완'에 있음\n"
                   "• 특히 회귀 팩터(30%)는 비중 대비 단독 성과가\n"
                   "  약함 → Trading Agent로 질적 보완 시 효과 극대화",
                   font_size=12, bg=BG_CARD)

    # ════════════════════════════════════════
    # 슬라이드 7: 기존전략 성과 요약
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "기존전략(A0) 성과 요약", "2021.01 ~ 2026.03 (63개월)")

    # 상단 KPI 카드 4개
    kpis = [
        ("누적수익률", "+140.1%", "KOSPI +107.3%"),
        ("CAGR", "18.2%", "KOSPI 14.9%"),
        ("Sharpe Ratio", "0.872", "KOSPI 0.703"),
        ("MDD", "-34.9%", "KOSPI -33.9%"),
    ]
    for i, (label, value, sub) in enumerate(kpis):
        left = 0.5 + i * 3.15
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(left), Inches(1.7), Inches(2.9), Inches(1.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        shape.line.width = Pt(0.5)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.font.name = "맑은 고딕"
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(28)
        p2.font.color.rgb = GREEN if not value.startswith("-") else RED
        p2.font.bold = True
        p2.font.name = "맑은 고딕"
        p3 = tf.add_paragraph()
        p3.text = sub
        p3.font.size = Pt(11)
        p3.font.color.rgb = LIGHT_GRAY
        p3.font.name = "맑은 고딕"

    # 추가 KPI
    sub_kpis = [
        ("월 승률", "59%"),
        ("월평균 수익률", "+1.59%"),
        ("연환산 변동성", "21.9%"),
        ("평균 회전율", "27.2%"),
    ]
    for i, (label, value) in enumerate(sub_kpis):
        left = 0.5 + i * 3.15
        add_text_box(slide, left, 3.15, 1.5, 0.3, label, font_size=11, color=GRAY)
        add_text_box(slide, left + 1.5, 3.15, 1.4, 0.3, value, font_size=11, bold=True, color=WHITE)

    # 연도별 비교 테이블
    add_text_box(slide, 0.5, 3.7, 12, 0.4, "연도별 성과 비교", font_size=16, bold=True, color=BLUE)
    add_table(slide, 0.5, 4.2, 12.3, [
        ["연도", "기존전략", "KOSPI200", "초과수익", "Sharpe", "MDD", "판정"],
        ["2021", "+19.5%", "-1.1%", "+20.6%p", "1.15", "-13.8%", "WIN"],
        ["2022", "-29.0%", "-26.7%", "-2.3%p", "-1.61", "-25.6%", "LOSE"],
        ["2023", "+19.5%", "+24.4%", "-4.9%p", "1.20", "-12.4%", "LOSE"],
        ["2024", "+6.0%", "-11.9%", "+17.9%p", "0.59", "-10.4%", "WIN"],
        ["2025", "+80.7%", "+96.4%", "-15.7%p", "2.38", "-5.0%", "LOSE"],
        ["2026(3M)", "+23.5%", "+32.8%", "-9.3%p", "2.92", "-5.2%", "LOSE"],
    ], col_widths=[1.5, 1.7, 1.7, 1.7, 1.5, 1.5, 1.2], font_size=12)

    # ════════════════════════════════════════
    # 슬라이드 8: 기존전략 시사점
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "기존전략 시사점", "강점과 약점 분석")

    # 왼쪽: 강점
    add_text_box(slide, 0.5, 1.6, 6, 0.4, "강점", font_size=20, bold=True, color=GREEN)

    add_accent_box(slide, 0.5, 2.2, 5.8, 1.0,
                   "하락장 방어력이 뛰어남\n"
                   "KOSPI 하락 월(27개월)에서 67%의 확률로 KOSPI보다 선방\n"
                   "하락장 평균: 전략 -2.57% vs KOSPI -3.97% → +1.40%p 초과",
                   font_size=12, bg=RGBColor(0xEC, 0xFD, 0xF5))

    add_accent_box(slide, 0.5, 3.4, 5.8, 0.9,
                   "롤링 12개월 기준 67% 확률로 KOSPI를 초과\n"
                   "평균 초과수익 +3.0%p, 최대 +19.8%p",
                   font_size=12, bg=RGBColor(0xEC, 0xFD, 0xF5))

    add_accent_box(slide, 0.5, 4.5, 5.8, 0.9,
                   "2021, 2024년 KOSPI 대비 대폭 초과 달성\n"
                   "2021: KOSPI -1.1% → 전략 +19.5% (+20.6%p)\n"
                   "2024: KOSPI -11.9% → 전략 +6.0% (+17.9%p)",
                   font_size=12, bg=RGBColor(0xEC, 0xFD, 0xF5))

    # 오른쪽: 약점/개선점
    add_text_box(slide, 7, 1.6, 6, 0.4, "약점 / 개선 포인트", font_size=20, bold=True, color=RED)

    add_accent_box(slide, 7, 2.2, 5.8, 1.0,
                   "상승장에서 KOSPI를 따라가지 못함\n"
                   "KOSPI 상승 월(36개월) 평균: 전략 +4.71% vs KOSPI +5.52%\n"
                   "→ 상승장 초과수익 -0.81%p (뒤처짐)",
                   font_size=12, bg=RGBColor(0xFE, 0xF2, 0xF2))

    add_accent_box(slide, 7, 3.4, 5.8, 0.9,
                   "2022년 하락장에서 KOSPI보다 더 빠짐\n"
                   "전략 -29.0% vs KOSPI -26.7% → -2.3%p 언더퍼폼\n"
                   "하락장에서 방어 실패한 해",
                   font_size=12, bg=RGBColor(0xFE, 0xF2, 0xF2))

    add_accent_box(slide, 7, 4.5, 5.8, 0.9,
                   "2025~2026 대세 상승장에서 KOSPI 대비 열위\n"
                   "2025: 전략 +80.7% vs KOSPI +96.4% → -15.7%p\n"
                   "2026: 전략 +23.5% vs KOSPI +38.7% → -15.1%p",
                   font_size=12, bg=RGBColor(0xFE, 0xF2, 0xF2))

    # 하단 종합
    add_accent_box(slide, 0.5, 5.8, 12.3, 1.2,
                   "종합 시사점\n\n"
                   "• 기존전략은 '하락장 방어 + 안정적 수익'에 강점 → 밸류 팩터의 본질적 특성\n"
                   "• 대세 상승장에서는 KOSPI를 따라가지 못함 → 모멘텀/성장 팩터 비중 조절 필요\n"
                   "• 2022년처럼 전방위 하락장에서는 밸류도 방어 불가 → Trading Agent의 레짐 판단이 필요한 이유",
                   font_size=13)

    # ════════════════════════════════════════
    # 슬라이드 9: 언더퍼폼 Deep Dive (표)
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "KOSPI 대비 언더퍼폼 원인 분석", "63개월 중 언더퍼폼 29개월(46%) — 최악 5개월 상세")

    # 표 1: 최악 5개월 요약
    add_text_box(slide, 0.5, 1.5, 12, 0.3, "최악 5개월 요약", font_size=15, bold=True, color=BLUE)
    add_table(slide, 0.5, 1.9, 12.3, [
        ["날짜", "전략", "KOSPI", "초과", "Top 섹터", "Top3 섹터 비중"],
        ["2025-09", "+5.5%", "+13.2%", "-7.7%p", "금융 27%", "62.3%"],
        ["2026-02", "+13.3%", "+19.3%", "-6.0%p", "전기·전자 28%", "55.4%"],
        ["2025-07", "-4.9%", "+1.1%", "-6.0%p", "금융 26%", "58.1%"],
        ["2024-03", "-2.1%", "+3.6%", "-5.7%p", "운송·창고 26%", "65.2%"],
        ["2026-03", "-5.0%", "+0.0%", "-5.0%p", "금융 23%", "54.3%"],
    ], col_widths=[1.5, 1.3, 1.3, 1.3, 3.0, 2.0], font_size=11)

    # 표 2: 반복 손실 종목
    add_text_box(slide, 0.5, 4.3, 12, 0.3, "반복 손실 종목 (밸류 트랩)", font_size=15, bold=True, color=RED)
    add_table(slide, 0.5, 4.7, 12.3, [
        ["종목", "손실 기여 등장", "섹터", "특징"],
        ["기아", "4회", "운송장비", "저PER (5~7배) — 자동차 업황 둔화 시 반복 하락"],
        ["HMM", "3회", "운송·창고", "저PBR (0.3~0.5배) — 해운 사이클 하락기 밸류 트랩"],
        ["현대글로비스", "3회", "운송·창고", "저PER — 물류 업황과 동반 하락"],
        ["한국전력", "3회", "전기·가스", "저PBR — 정책 리스크, 만년 저평가"],
        ["SK스퀘어", "2회", "금융", "저PER — SK하이닉스 연동 변동성"],
    ], col_widths=[1.8, 1.5, 1.5, 5.5], font_size=11)

    # ════════════════════════════════════════
    # 슬라이드 9-2: 언더퍼폼 원인 요약
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "언더퍼폼 원인 요약", "두 가지 핵심 문제")

    # ① 섹터 쏠림
    add_text_box(slide, 0.5, 1.6, 12, 0.4, "① 섹터 쏠림", font_size=20, bold=True, color=RED)

    add_accent_box(slide, 0.5, 2.1, 12.3, 1.6,
                   "'싼 주식'을 고르다 보니 비슷한 업종이 몰린다\n\n"
                   "• 최악 10개월 평균: Top 3 섹터 합계 59% — 30종목이지만 실질 3~4개 업종에 집중\n"
                   "• 금융(10회), 운송·창고(9회), 운송장비(8회) 반복 등장\n"
                   "• 이 업종이 동시에 빠지면 포트폴리오 전체가 끌려감",
                   font_size=13, bg=RGBColor(0xFE, 0xF2, 0xF2))

    # ② 밸류 트랩
    add_text_box(slide, 0.5, 4.0, 12, 0.4, "② 밸류 트랩", font_size=20, bold=True, color=RED)

    add_accent_box(slide, 0.5, 4.5, 12.3, 1.6,
                   "팩터 점수가 높아서 매번 뽑히지만, 주가는 계속 빠지는 종목\n\n"
                   "• 기아(4회), HMM(3회), 현대글로비스(3회), 한국전력(3회) — 손실 기여 상위 단골\n"
                   "• 저PER·저PBR이라 모델은 '저평가'로 보지만, 시장은 '이유가 있어서 싸다'고 판단\n"
                   "• 단, 장기적으로는 수익을 내는 종목이라 단순 제거 시 오히려 성과 악화",
                   font_size=13, bg=RGBColor(0xFE, 0xF2, 0xF2))

    # 한줄 요약
    add_accent_box(slide, 0.5, 6.4, 12.3, 0.7,
                   "한줄 요약: 모델이 '싸다'에 과집중 → 같은 업종·같은 종목 반복 → 빠질 때 같이 빠진다",
                   font_size=14, bg=BG_CARD)

    # ════════════════════════════════════════
    # 슬라이드 10: 섹터 캡 실험
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "개선 실험 ①: 섹터 종목 수 상한", "업종 쏠림을 제한하면?")

    add_table(slide, 1, 1.8, 11, [
        ["전략", "설명", "수익률", "CAGR", "Sharpe", "MDD"],
        ["A0 기준", "섹터 제한 없음", "+149.2%", "+19.0%", "0.90", "-34.3%"],
        ["섹터 상한 25%", "한 업종 최대 7~8종목", "+158.2%", "+19.8%", "0.93", "-34.2%"],
        ["섹터 상한 20%", "한 업종 최대 6종목", "+150.4%", "+19.1%", "0.92", "-34.7%"],
        ["섹터 상한 15%", "한 업종 최대 4~5종목", "+113.5%", "+15.5%", "0.78", "-37.7%"],
    ], col_widths=[2.2, 2.8, 1.3, 1.2, 1.2, 1.2], font_size=13)

    add_accent_box(slide, 0.5, 4.5, 5.8, 2.0,
                   "결과: Sharpe 소폭 개선, MDD 변화 없음\n\n"
                   "• 섹터 상한 25%: Sharpe 0.90→0.93\n"
                   "  수익률도 149%→158%로 향상\n"
                   "• 그러나 MDD는 -34.3%→-34.2%로 거의 동일\n"
                   "• 15%로 너무 조이면 오히려 성과 악화",
                   font_size=12, bg=RGBColor(0xEC, 0xFD, 0xF5))

    add_accent_box(slide, 6.8, 4.5, 5.8, 2.0,
                   "시사점\n\n"
                   "• 섹터 분산은 수익 효율 개선에는 유효\n"
                   "• 하지만 MDD 문제는 해결되지 않음\n"
                   "• 업종 쏠림은 원인의 일부일 뿐\n"
                   "  → 밸류 트랩 문제가 남아있음",
                   font_size=12, bg=BG_CARD)

    # ════════════════════════════════════════
    # 슬라이드 10-2: 밸류 트랩 필터 실험
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "개선 실험 ①: 밸류 트랩 필터", "3개월 연속 하락 종목을 제외하면?")

    add_table(slide, 1, 1.8, 11, [
        ["전략", "설명", "수익률", "CAGR", "Sharpe", "MDD"],
        ["A0 기준", "현재 전략", "+149.2%", "+19.0%", "0.90", "-34.3%"],
        ["밸류트랩 필터", "3개월 연속 하락 종목 제외", "+101.4%", "+14.3%", "0.71", "-37.7%"],
    ], col_widths=[2.2, 3.5, 1.3, 1.2, 1.2, 1.2], font_size=13)

    add_accent_box(slide, 0.5, 3.5, 5.8, 2.0,
                   "결과: 오히려 악화\n\n"
                   "• 수익률 149% → 101% (수익 1/3 날아감)\n"
                   "• Sharpe 0.90 → 0.71, MDD도 악화\n"
                   "• 월평균 3.8종목이 필터에 걸림\n"
                   "• 기아·HMM 등은 일시 하락 후 반등하는 종목\n"
                   "  → 제거하면 반등 구간도 놓침",
                   font_size=12, bg=RGBColor(0xFE, 0xF2, 0xF2))

    add_accent_box(slide, 6.8, 3.5, 5.8, 2.0,
                   "시사점\n\n"
                   "• 밸류 트랩은 밸류 전략의 본질적 리스크\n"
                   "• 종목을 빼는 건 답이 아님 (반등도 놓침)\n"
                   "• 핵심은 '어떤 종목을 빼느냐'가 아니라\n"
                   "  '언제 포지션을 줄이느냐' → 시점 관리\n"
                   "• → 다음 슬라이드: 시장 레짐 필터",
                   font_size=12, bg=BG_CARD)

    add_accent_box(slide, 0.5, 5.8, 12.3, 1.0,
                   "밸류 트랩의 해결책은 '종목 제거'가 아니라 '시점 관리'\n"
                   "→ 밸류 트랩 종목이 빠지는 시기(하락장)에 포지션 자체를 줄이면 피해를 줄일 수 있다",
                   font_size=13)

    # ════════════════════════════════════════
    # 슬라이드 10-2: 시장 레짐 필터 실험
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "개선 실험 ②: 시장 레짐 필터", "KOSPI 12개월 이동평균 하회 시 현금 비중 확대")

    add_table(slide, 0.8, 1.8, 11.5, [
        ["전략", "수익률", "CAGR", "Sharpe", "MDD"],
        ["A0 기준 (현금 0%)", "+149.2%", "+19.0%", "0.90", "-34.3%"],
        ["하락장 현금 30%", "+159.9%", "+20.0%", "1.00", "-26.9%"],
        ["하락장 현금 50%", "+166.2%", "+20.5%", "1.06", "-21.7%"],
        ["하락장 현금 70%", "+171.8%", "+21.0%", "1.12", "-16.3%"],
    ], col_widths=[3.5, 1.8, 1.5, 1.5, 1.5], font_size=13)

    add_accent_box(slide, 0.5, 4.3, 5.8, 1.5,
                   "결과: 모든 지표가 개선\n\n"
                   "• MDD: -34.3% → -16.3% (현금 70%)\n"
                   "• Sharpe: 0.90 → 1.12\n"
                   "• 수익률도 149% → 172%로 오히려 상승\n"
                   "  (하락장에서 덜 잃으니 복리 효과)",
                   font_size=12, bg=RGBColor(0xEC, 0xFD, 0xF5))

    add_accent_box(slide, 6.8, 4.3, 5.8, 1.5,
                   "원리\n\n"
                   "• KOSPI가 12개월 이동평균 하회 = 하락장 신호\n"
                   "• 63개월 중 26개월(41%)이 하락장 (주로 2022년)\n"
                   "• 이 기간에 현금 비중을 높여 손실을 줄임\n"
                   "• 단순한 지표 하나로 MDD 절반 축소",
                   font_size=12, bg=BG_CARD)

    add_accent_box(slide, 0.5, 6.1, 12.3, 0.8,
                   "밸류 트랩 문제도 해결: 종목을 빼는 게 아니라 '빠지는 시기'에 포지션을 줄여 피해 축소\n"
                   "주의: 2022년 하락장에 집중된 효과 — 다양한 시장 환경에서의 추가 검증 필요",
                   font_size=12)

    # ════════════════════════════════════════
    # 슬라이드 10-4: 조합 결과
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "개선 조합: 섹터 캡 + 레짐 필터", "두 가지를 합치면?")

    add_table(slide, 0.5, 1.8, 12.3, [
        ["전략", "수익률", "CAGR", "Sharpe", "MDD"],
        ["A0 기준 (현재)", "+149.2%", "+19.0%", "0.90", "-34.3%"],
        ["섹터 캡 25% only", "+158.2%", "+19.8%", "0.93", "-34.2%"],
        ["레짐 필터 50% only", "+166.2%", "+20.5%", "1.06", "-21.7%"],
        ["섹터25% + 레짐50%", "+173.5%", "+21.1%", "1.09", "-21.7%"],
        ["레짐 필터 70% only", "+171.8%", "+21.0%", "1.12", "-16.3%"],
        ["섹터25% + 레짐70%", "+178.2%", "+21.5%", "1.14", "-16.3%"],
    ], col_widths=[3.5, 2.0, 1.8, 1.5, 1.5], font_size=12)

    add_accent_box(slide, 0.5, 5.0, 5.8, 1.5,
                   "조합 효과\n\n"
                   "• 섹터 캡: 수익률을 +7~8%p 추가 향상\n"
                   "• 레짐 필터: MDD를 절반으로 축소\n"
                   "• 두 가지는 독립적으로 작동 → 합치면 효과 합산\n"
                   "• 최고 조합: 섹터25% + 레짐70%",
                   font_size=12, bg=RGBColor(0xEC, 0xFD, 0xF5))

    add_accent_box(slide, 6.8, 5.0, 5.8, 1.5,
                   "A0 대비 개선 폭\n\n"
                   "• 수익률: 149% → 178% (+29%p)\n"
                   "• Sharpe: 0.90 → 1.14 (+27%)\n"
                   "• MDD: -34.3% → -16.3% (절반 이하)\n"
                   "• 단순 규칙 두 개만으로 대폭 개선",
                   font_size=12, bg=BG_CARD)

    add_accent_box(slide, 0.5, 6.7, 12.3, 0.5,
                   "향후: 레짐 판단 고도화 (AI Agent), 현금 비중 동적 조절, 더 긴 백테스트 기간으로 추가 검증",
                   font_size=12)

    # ════════════════════════════════════════
    # 슬라이드 11: 왜 AI Trading Agent인가
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "다음 단계: AI Trading Agent", "퀀트 모델이 못 보는 '비정형 데이터'를 AI가 읽는다")

    # ── 좌측: 문제 → 왜 AI가 필요한가 ──
    add_text_box(slide, 0.5, 1.5, 6, 0.3, "퀀트 모델의 한계: 정형 데이터만 본다", font_size=15, bold=True, color=RED)

    add_accent_box(slide, 0.5, 2.0, 5.8, 1.8,
                   "현재 모델이 사용하는 데이터 (정형)\n\n"
                   "  PER, PBR, ROE, 매출성장률, EPS 컨센서스...\n"
                   "  → 숫자로 표현되는 재무 데이터만 활용\n\n"
                   "현재 모델이 못 보는 데이터 (비정형)\n\n"
                   "  뉴스: '기아 미국 관세 리스크 확대'\n"
                   "  공시: '한국전력 전기요금 동결 결정'\n"
                   "  시장심리: '해운주 투매 심리 확산'\n"
                   "  거시경제: '미국 금리 인하 기대감 후퇴'",
                   font_size=11, bg=RGBColor(0xFE, 0xF2, 0xF2))

    add_accent_box(slide, 0.5, 4.0, 5.8, 0.8,
                   "앞서 확인한 문제들이 바로 이 한계에서 발생\n"
                   "• 밸류 트랩: 숫자는 '싸다'지만 뉴스는 '이유가 있다'\n"
                   "• 레짐 전환: 12개월 이평은 후행 — 뉴스가 더 빠르다",
                   font_size=11, bg=BG_CARD)

    # ── 우측: TradingAgents 구조 ──
    add_text_box(slide, 7, 1.5, 6, 0.3, "TradingAgents: 비정형 데이터 분석 AI", font_size=15, bold=True, color=BLUE)

    # 파이프라인 — 세로 플로우
    add_accent_box(slide, 7, 2.0, 5.8, 1.1,
                   "애널리스트 팀 (4개 AI 에이전트가 독립 분석)\n\n"
                   "  펀더멘탈   재무제표 기반 내재가치  |  뉴스   글로벌 뉴스·거시경제\n"
                   "  센티먼트   소셜미디어·시장심리     |  기술   MACD, RSI 등 차트",
                   font_size=10, bg=RGBColor(0xEF, 0xF6, 0xFF))

    add_text_box(slide, 9.5, 3.15, 1, 0.25, "▼", font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_accent_box(slide, 7, 3.4, 5.8, 0.7,
                   "리서처 토론 — Bullish vs Bearish 교차 검증\n"
                   "강세/약세 관점으로 토론 → 편향 제거, 균형 잡힌 결론",
                   font_size=10, bg=RGBColor(0xFE, 0xF9, 0xEF))

    add_text_box(slide, 9.5, 4.15, 1, 0.25, "▼", font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_accent_box(slide, 7, 4.4, 5.8, 0.7,
                   "트레이더 → 리스크 관리 → 포트폴리오 매니저\n"
                   "매매 제안 → 위험 평가 → 최종 승인/거절",
                   font_size=10, bg=RGBColor(0xEC, 0xFD, 0xF5))

    # ── 하단: 이현 전략 적용 포인트 ──
    add_text_box(slide, 0.5, 5.3, 12, 0.3, "이현 퀀트 전략에 적용", font_size=15, bold=True, color=BLUE)

    add_accent_box(slide, 0.5, 5.7, 3.9, 1.2,
                   "밸류 트랩 판단\n\n"
                   "뉴스+펀더멘탈 에이전트가\n"
                   "'업황 둔화' vs '일시 조정' 구분\n"
                   "→ 규칙 실패 영역을 AI가 커버",
                   font_size=10, bg=BG_CARD)

    add_accent_box(slide, 4.7, 5.7, 3.9, 1.2,
                   "레짐 전환 조기 감지\n\n"
                   "센티먼트+뉴스 에이전트가\n"
                   "이동평균보다 빠르게 하락장 감지\n"
                   "→ 현금 비중 동적 결정",
                   font_size=10, bg=BG_CARD)

    add_accent_box(slide, 8.9, 5.7, 3.9, 1.2,
                   "섹터 가중치 동적 조절\n\n"
                   "리서처 토론으로 업종 전망 판단\n"
                   "'반도체 호조' → 상한 완화\n"
                   "'해운 하락' → 상한 강화",
                   font_size=10, bg=BG_CARD)

    # ════════════════════════════════════════
    # 슬라이드 12: 지표 설명
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "주요 지표 설명", "Q&A")

    metrics = [
        ("MDD", "백테스트 기간 중 고점 대비 가장 많이 빠진 비율\n매 시점에서 역대 최고점 대비 하락률 → 그 중 최대값"),
        ("평균 회전율", "매월 (신규 편입 + 편출) / (2 × 전체 종목 수)\n높으면 거래비용 증가 → 백테스트에서 편도 30bp 반영"),
        ("HHI", "포트폴리오 집중도: 각 종목 비중(%)² 합산\n30종목 균등이면 ≈333 / 높을수록 소수 종목에 집중"),
        ("승률 (롤링 24개월)", "24개월 윈도우 이동하며 KOSPI200 대비 초과수익\n양수인 구간의 비율 (전략 간 비교가 아닌 vs 벤치마크)"),
        ("IS/OOS", "과적합 검증 방법\nIS: 학습 기간 (좋은 게 당연) / OOS: 못 본 기간 (여기서도 좋아야 진짜)"),
    ]

    for i, (name, desc) in enumerate(metrics):
        row = i // 2
        col = i % 2
        left = 0.5 + col * 6.4
        top = 1.7 + row * 1.7
        add_text_box(slide, left, top, 1.5, 0.35, name, font_size=14, bold=True, color=BLUE)
        add_text_box(slide, left + 1.6, top, 4.5, 1.3, desc, font_size=12, color=WHITE)

    # ════════════════════════════════════════
    # 슬라이드 15: 수정 반영
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "수정 반영 사항")

    add_table(slide, 2, 1.8, 9, [
        ["#", "요청", "반영"],
        ["1", "회귀only 전략 제거", "✅ 완료"],
        ["2", "성과지표 월별/연도별 집계", "✅ 완료"],
        ["3", "최대 비중 15% → 10%", "✅ 완료"],
        ["4", "시총 하한 1조 → 5천억", "✅ 완료"],
        ["5", "히트맵에 수치 표시", "✅ 완료"],
        ["6", "포트폴리오 특성 단순평균 표시", "✅ 완료"],
        ["7", "유니버스 KOSPI만", "✅ 완료"],
    ], col_widths=[0.8, 5.0, 3.2], font_size=14)

    # ════════════════════════════════════════
    # 슬라이드 16: 로드맵
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    slide_title(slide, "추후 작업 로드맵")

    add_table(slide, 1, 1.8, 11, [
        ["작업", "내용"],
        ["기간 실험", "5년 vs 3년 vs 2년 백테스트 비교"],
        ["비중 캡 실험", "10% vs 8% vs 7%"],
        ["팩터 비중 시뮬레이션", "전략실험실에서 다양한 조합 테스트"],
        ["Trading Agent 개발", "회귀 팩터 보완 (가치함정 필터, 모멘텀 구분)"],
        ["이현 펀드 데이터 비교", "실제 운용 성과 vs 백테스트 성과 분석"],
    ], col_widths=[3.5, 7.5], font_size=14)

    add_text_box(slide, 1, 5.5, 11, 0.4, "미팅 시 이현 실무진에게 요청할 데이터",
                 font_size=16, bold=True, color=YELLOW)
    add_bullet_list(slide, 1.2, 6.0, 10, 1.5, [
        "1. 펀드 벤치마크 월별 수익률",
        "2. 펀드 기준가 시계열 (일별 또는 월별)",
        "3. 무위험수익률 (CD91일 등)",
        "4. 실제 운용 회전율 데이터",
    ], font_size=13)

    # ─── 저장 ───
    out_path = "/Users/namsugyeong/Desktop/alpha_lab/analysis/퀀트모델_리뷰_회귀팩터분석.pptx"
    prs.save(out_path)
    print(f"PPT 저장 완료: {out_path}")


if __name__ == "__main__":
    main()