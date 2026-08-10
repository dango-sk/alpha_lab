# -*- coding: utf-8 -*-
"""레짐 모델 정리 PPT (4장) — 쉬운 말 + 실험 표 2개 + 차트. 박사님 보고용."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

CH = Path("analysis/ppt_charts")
NAVY=RGBColor(0x1F,0x33,0x55); BLUEc=RGBColor(0x2E,0x5B,0x9A); GR=RGBColor(0x44,0x44,0x44)
REDc=RGBColor(0xB1,0x2A,0x2A); GRNc=RGBColor(0x1B,0x6E,0x3C); HDR=RGBColor(0x2E,0x5B,0x9A); HL=RGBColor(0xE8,0xF1,0xE8)
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)

def slide(title):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    tf=s.shapes.add_textbox(Inches(0.45),Inches(0.28),Inches(12.4),Inches(0.85)).text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=title; p.font.size=Pt(25); p.font.bold=True; p.font.color.rgb=NAVY
    return s

def body(s, items, left, top, width, size=13.5):
    tf=s.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(5.6)).text_frame; tf.word_wrap=True
    for i,(lvl,txt,col,bold) in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=txt; p.level=lvl; p.font.size=Pt(size-lvl*1.5); p.font.bold=bold; p.font.color.rgb=col; p.space_after=Pt(4)

def table(s, data, left, top, width, height, hi_row=None, fs=11.5, col_w=None):
    rows,cols=len(data),len(data[0])
    gf=s.shapes.add_table(rows,cols,Inches(left),Inches(top),Inches(width),Inches(height)); t=gf.table
    if col_w:
        for j,w in enumerate(col_w): t.columns[j].width=Inches(w)
    for i,row in enumerate(data):
        for j,val in enumerate(row):
            c=t.cell(i,j); c.text=str(val)
            pr=c.text_frame.paragraphs[0]; pr.font.size=Pt(fs); pr.alignment=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT
            if i==0:
                pr.font.bold=True; pr.font.color.rgb=RGBColor(0xFF,0xFF,0xFF); c.fill.solid(); c.fill.fore_color.rgb=HDR
            else:
                c.fill.solid(); c.fill.fore_color.rgb=(HL if hi_row==i else RGBColor(0xFF,0xFF,0xFF))
                pr.font.bold=(hi_row==i)
    return t

# ===== Slide 1: 왜 HSMM =====
s=slide("① 모델 구조 선택 — '지속학습(HSMM)'을 고른 이유")
body(s,[
 (0,"약세장 신호가 너무 자주 켜졌다 꺼졌다 해서, 안정화 방법 5가지를 같은 조건에서 비교했습니다.",GR,False),
],0.5,1.15,12.2,13)
table(s,[
 ["안정화 방법","위기탐지(7중)","선행(개월)","헛경보","들락날락","약세장 구분점수"],
 ["기본 (그냥 HMM)","7","+0.0","5","35","+1.65"],
 ["끈끈하게(Sticky)","7","+0.4","5","23","-0.70 (구분 붕괴)"],
 ["전환에 벌점(Penalty)","7","-1.1","4","23","+1.40"],
 ["문턱 차등(Hysteresis)","7","-1.0","5","31","+1.67"],
 ["지속학습(HSMM) ★채택","7","+0.0","3","31","+1.75 (최고)"],
], 0.5, 1.9, 12.3, 2.6, hi_row=5, col_w=[3.3,2.0,1.8,1.6,1.8,1.8])
body(s,[
 (0,"→ HSMM은 위기탐지(7/7)·선행을 그대로 유지하면서, 약세장 구분점수가 가장 높고(+1.75) 헛경보가 가장 적었습니다(3건).",GRNc,True),
 (0,"   (\"약세장 구분점수\" = 약세라고 한 달과 강세라고 한 달의 실제 수익 차이. 높을수록 잘 가려냄)",GR,False),
 (0,"   참고: 같은 HSMM도 '지속기간 학습'을 끄면 구분점수가 +1.5 → +0.8로 반토막 → 지속학습이 핵심.",GR,False),
],0.5,4.7,12.3,12.5)

# ===== Slide 2: 최종 팩터 =====
s=slide("② 최종 선택한 정보(팩터) — 6개, 그리고 뺀 것들")
body(s,[(0,"여러 정보를 넣고 빼며 검증한 결과입니다.",GR,False)],0.5,1.1,12,13)
# 후보 표
table(s,[
 ["추가 시도한 정보","약세장 구분","모델 안정성","실제 수익(FCF)","채택?"],
 ["환율(원화 3개월 변화)","유지","크게 좋아짐 ★","—","✅ 채택"],
 ["금리(미국 10년)","좋아짐(+2.3)","안정","도움 안 됨","❌"],
 ["달러(DXY)","조금 좋아짐","안정","—","❌"],
 ["뉴스(AI 점수)","설명력 있음","망가짐","—","❌"],
 ["구리/금 비율","도움 안 됨","—","—","❌"],
], 0.5, 1.7, 7.7, 2.7, hi_row=1, fs=11, col_w=[2.6,1.4,1.5,1.3,0.9])
body(s,[
 (0,"■ 최종 모델에 들어간 6개",NAVY,True),
 (1,"① 오르는 종목 비율",GR,False),
 (1,"② 그 비율의 변화",GR,False),
 (1,"③ 52주 신저가 종목 비율",GR,False),
 (1,"④ 그 비율의 변화",GR,False),
 (1,"⑤ 모멘텀(1개월−6개월)",GR,False),
 (1,"⑥ 환율(원화 3개월 변화) ★",GRNc,True),
 (0,"= 시장 내부 5개 + 자본흐름(환율) 1개",BLUEc,True),
], 8.5, 1.7, 4.5, 13.5)
body(s,[
 (0,"→ '정보를 더 넣으면 좋아진다'가 아니었습니다. 환율만 빼고 다 기각 — 환율은 예측을 잘해서가 아니라 결과를 안 흔들리게 해줘서 채택.",GRNc,True),
],0.5,4.8,12.3,12.5)

# ===== Slide 3: 환율 운빨 =====
s=slide("③ '환율'이 모델을 운빨에서 구해준 그림")
body(s,[
 (0,"시장 내부지표만 쓰면 — 컴퓨터 랜덤 초기값에 따라 결과가 +/− 로 뒤집혔습니다 (운빨).",REDc,False),
 (0,"환율을 넣으니 — 어떤 초기값으로 돌려도 항상 약세장을 잘 구분했습니다.",GRNc,True),
 (0,"오른쪽: 5가지 초기값으로 돌린 '약세장 구분 점수'",GR,False),
 (1,"빨강(내부지표만): 들쭉날쭉, 마이너스도 나옴",REDc,False),
 (1,"초록(+환율): 5개 모두 비슷하게 +",GRNc,False),
 (0,"→ 그래서 환율의 진짜 역할은 '예측력'이 아니라 '안정성'입니다.",NAVY,True),
],0.5,1.5,6.0,14)
s.shapes.add_picture(str(CH/"k1.png"), Inches(6.8), Inches(1.8), width=Inches(6.2))

# ===== Slide 4: 결과 + overlay 돌파구 =====
s=slide("④ 결과: 위기는 우리가 잘 잡고, 수익은 '비중 줄이기'로 따라잡음")
body(s,[
 (0,"우리 모델 vs 기존 AI: 위기탐지·낙폭방어는 우위, 단 그냥 쓰면 수익은 AI가 높았음.",GR,False),
 (0,"이유: 2022 하락장에 우리가 잠깐 반등에 '강세'로 깜빡임 + 반등장에 손절로 바닥 매도.",GR,False),
 (0,"→ 해결: 약세장에 전략을 통째로 바꾸지 말고 '비중만 줄이기'(강세전략 유지)",NAVY,True),
],0.5,1.15,12.4,13)
table(s,[
 ["방식","누적수익","Sharpe","최대낙폭"],
 ["그냥 강세전략만","368%","0.78","-43.7%"],
 ["약세장에 전략 교체 (기존 시도)","353%","0.75","-38.1%"],
 ["약세장에 비중만 50% 축소 ★신규","387%","0.82","-37.2%"],
 ["기존 AI(ai_v2)","406%","0.82","-40.6%"],
], 0.5, 2.5, 8.3, 2.2, hi_row=3, fs=12, col_w=[3.6,1.7,1.4,1.6])
body(s,[
 (0,"→ '비중만 줄이기'가 그냥 강세전략보다 수익↑·낙폭↓,",GRNc,True),
 (0,"   기존 AI와 Sharpe 동률(0.82)·낙폭은 더 우수(-37 vs -41).",GRNc,True),
 (0,"(거래비용 반영치. AI의 높은 수익은 '뉴스 정보' 덕이고, 우리는 그 정보 없이 비중조절로 근접)",GR,False),
],0.5,4.9,12.3,13)

out=Path("레짐모델_정리_표포함.pptx"); prs.save(out)
print(f"저장: {out.resolve()} (4장)")
