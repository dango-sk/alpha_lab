# -*- coding: utf-8 -*-
"""
analysis/make_hsmm_ppt_pm.py
운용역(PM)용 요약 덱 — '낙폭 방어·수익 보존·주식비중' 언어.
반영: 이현 스크리닝 전략(FCF아님) 명칭 / KOSPI 시장 비교 / 미투자현금 연2.5% 운용 / vol-targeting=꼬리보험.
출력: HSMM_레짐_운용역_보고.pptx
"""
import numpy as np, pandas as pd
from pathlib import Path
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
plt.rcParams["font.family"] = "AppleGothic"; plt.rcParams["axes.unicode_minus"] = False
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = Path(__file__).parent.parent
A = Path(__file__).parent
ASSET = A / "ppt_assets_pm"; ASSET.mkdir(exist_ok=True)
OUT = BASE / "HSMM_레짐_운용역_보고.pptx"
FONT = "Apple SD Gothic Neo"
CASH = 0.025

NAVY = RGBColor(0x1F, 0x38, 0x64); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x26, 0x2A, 0x33); GRAY = RGBColor(0x59, 0x5F, 0x6B)
BLUE = RGBColor(0x2F, 0x54, 0x96); GREEN = RGBColor(0x2E, 0x74, 0x4F)
RED = RGBColor(0xC0, 0x3A, 0x3A); ALT = RGBColor(0xF2, 0xF4, 0xF7); HLG = RGBColor(0xE2, 0xEF, 0xDA)
CBM, COV, CMK = "#C03A3A", "#1F3864", "#8A93A6"

# ───────── 데이터 & 지표 ─────────
df = pd.read_csv(A / "fcf_overlay_series.csv")
corr = np.corrcoef(df.expA, df.expP)[0, 1]
cm = (1 + CASH) ** (1 / 12) - 1
ov_cash = df.ovP.values + (1 - df.expP.values) * cm          # 리스크 오버레이 + 현금2.5%
ovA_cash = df.ovA.values + (1 - df.expA.values) * cm
yrs = df.ym.str.slice(0, 4).values
ticks = [i for i in range(len(df)) if i == 0 or yrs[i] != yrs[i - 1]]


def perf(r):
    r = np.asarray(r); eq = np.cumprod(1 + r); y = len(r) / 12
    return eq[-1] ** (1 / y) - 1, r.mean() * 12 / (r.std() * np.sqrt(12)), (eq / np.maximum.accumulate(eq) - 1).min()


mk = perf(df.kospi.values); bm = perf(df.bench.values); ov = perf(ov_cash); va = perf(ovA_cash)

# 1) 누적수익 곡선 (KOSPI / 이현 / +오버레이)
fig, ax = plt.subplots(figsize=(6.2, 3.3), dpi=200)
x = np.arange(len(df))
ax.plot(x, np.cumprod(1 + df.kospi.values), color=CMK, lw=1.6, label="KOSPI 지수 (시장)")
ax.plot(x, np.cumprod(1 + df.bench.values), color=CBM, lw=2.0, label="이현 스크리닝 (단독)")
ax.plot(x, np.cumprod(1 + ov_cash), color=COV, lw=2.2, label="+ 리스크 오버레이 (채택)")
ax.set_xticks(ticks); ax.set_xticklabels([yrs[i] for i in ticks], fontsize=8)
ax.set_ylabel("누적 성장 (1 → 배수)")
ax.set_title("누적 성과 — 오버레이: 수익은 더 높게, 하락은 더 얕게", fontsize=11, pad=8)
ax.legend(fontsize=8.5, loc="upper left"); ax.spines[["top", "right"]].set_visible(False)
ax.axvspan(23, 24, color="#C03A3A", alpha=0.08); ax.axvspan(44, 57, color="#C03A3A", alpha=0.08)
fig.tight_layout(); fig.savefig(ASSET / "equity.png", bbox_inches="tight"); plt.close(fig)

# 2) 최대낙폭 비교 (KOSPI / 이현 / +오버레이 / +vol-tgt)
fig, ax = plt.subplots(figsize=(5.4, 3.3), dpi=200)
names = ["KOSPI\n(시장)", "이현\n스크리닝", "+ 리스크\n오버레이", "+ vol-\ntargeting"]
mdd = [mk[2] * 100, bm[2] * 100, ov[2] * 100, va[2] * 100]; cols = [CMK, CBM, COV, "#B7C0D0"]
bars = ax.bar(names, mdd, color=cols, width=0.62)
for b, v in zip(bars, mdd):
    ax.text(b.get_x() + b.get_width() / 2, v - 1.6, f"{v:.0f}%", ha="center", va="top", color="white", fontsize=10.5, fontweight="bold")
ax.axhline(0, color="#888", lw=0.8); ax.set_ylim(-50, 2); ax.set_ylabel("최대낙폭 MDD (%)")
ax.set_title("최대낙폭: 이현 단독 -44% → 오버레이 -33% (시장 수준)", fontsize=10.5, pad=8)
ax.spines[["top", "right"]].set_visible(False)
fig.tight_layout(); fig.savefig(ASSET / "mdd.png", bbox_inches="tight"); plt.close(fig)

# 3) 주식비중 겹침
fig, ax = plt.subplots(figsize=(6.0, 3.3), dpi=200)
ax.plot(x, df.expP, color=COV, lw=1.8, label="리스크 오버레이")
ax.plot(x, df.expA, color="#2E744F", lw=1.4, ls="--", label="+ vol-targeting")
ax.set_xticks(ticks); ax.set_xticklabels([yrs[i] for i in ticks], fontsize=8)
ax.set_ylabel("주식 비중"); ax.set_ylim(0.1, 1.05)
ax.set_title(f"두 방식의 주식비중이 거의 동일 → 약 99% 겹침", fontsize=11, pad=8)
ax.legend(fontsize=9, loc="lower right"); ax.spines[["top", "right"]].set_visible(False)
fig.tight_layout(); fig.savefig(ASSET / "exp.png", bbox_inches="tight"); plt.close(fig)

# ───────── PPT 헬퍼 ─────────
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]; SW, SH = prs.slide_width, prs.slide_height


def slide(): return prs.slides.add_slide(BLANK)


def tb(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    b = s.shapes.add_textbox(l, t, w, h); tf = b.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05); tf.margin_top = tf.margin_bottom = Inches(0.02)
    return b, tf


def run(p, text, size, color=INK, bold=False, italic=False):
    r = p.add_run(); r.text = text; f = r.font
    f.size = Pt(size); f.name = FONT; f.bold = bold; f.italic = italic; f.color.rgb = color
    return r


def header(s, num, title, subtitle=None):
    _, tf = tb(s, Inches(0.55), Inches(0.3), Inches(12.3), Inches(1.0))
    p = tf.paragraphs[0]; run(p, f"{num}  ", 28, BLUE, True); run(p, title, 28, NAVY, True)
    if subtitle:
        p2 = tf.add_paragraph(); run(p2, subtitle, 14, GRAY, italic=True)
    ln = s.shapes.add_shape(1, Inches(0.55), Inches(1.28), Inches(12.25), Pt(2.2))
    ln.fill.solid(); ln.fill.fore_color.rgb = NAVY; ln.line.fill.background()


def bullets(s, items, l, t, w, h, size=14, gap=10):
    _, tf = tb(s, l, t, w, h)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(gap)
        if isinstance(it, tuple):
            mark, txt, col, bold = it
            run(p, mark + "  ", size, col, True); run(p, txt, size, INK, bold)
        else:
            run(p, "•  ", size, NAVY, True); run(p, it, size, INK)


def table(s, rows, l, t, w, colw, hl=None, header_h=0.5, row_h=0.62, fs=13):
    nr, nc = len(rows), len(rows[0])
    gr = s.shapes.add_table(nr, nc, l, t, w, Inches(header_h + row_h * (nr - 1))).table
    gr.first_row = False; gr.horz_banding = False
    tot = sum(colw)
    for j, cw in enumerate(colw): gr.columns[j].width = Emu(int(w * cw / tot))
    gr.rows[0].height = Inches(header_h)
    for i in range(1, nr): gr.rows[i].height = Inches(row_h)
    for i, rowd in enumerate(rows):
        for j, val in enumerate(rowd):
            c = gr.cell(i, j); c.margin_left = c.margin_right = Inches(0.08)
            c.margin_top = c.margin_bottom = Inches(0.03); c.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i == 0: c.fill.solid(); c.fill.fore_color.rgb = NAVY
            elif hl is not None and i == hl: c.fill.solid(); c.fill.fore_color.rgb = HLG
            else: c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else ALT
            tfc = c.text_frame; tfc.word_wrap = True
            p = tfc.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            col = WHITE if i == 0 else (GREEN if (hl and i == hl) else INK)
            run(p, str(val), fs, col, bold=(i == 0 or (hl and i == hl)))


def pct(x): return f"{x*100:.1f}%"

# ───────── 슬라이드 ─────────
# 0. 표지
s = slide()
bar = s.shapes.add_shape(1, 0, Inches(2.3), SW, Inches(2.6)); bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
_, tf = tb(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(2.1), MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]; run(p, "레짐 기반 하방 방어 오버레이", 34, WHITE, True)
p2 = tf.add_paragraph(); run(p2, "이현 스크리닝 전략 적용 결과 — 낙폭은 시장 수준으로, 수익은 더 높게 (운용역 요약본)", 16, RGBColor(0xCF, 0xDA, 0xEC))
_, tf2 = tb(s, Inches(0.9), Inches(5.15), Inches(11.5), Inches(0.6))
run(tf2.paragraphs[0], "검증기간 2018.04 ~ 2026.07 (약 8년) · 미투자 현금 연 2.5% 운용 가정", 13, GRAY)

# 1. 한 장 요약
s = slide(); header(s, "①", "한 장 요약", "종목은 그대로 두고, 시장이 위험할 때 '주식 비중'만 자동으로 줄이는 장치")
bullets(s, [("무엇을", "시장 위험을 읽어 주식 비중을 조절하는 '리스크 스위치'를 이현 스크리닝 전략 위에 씌움 (종목선정은 그대로)", NAVY, True),
            ("최대낙폭", f"{pct(bm[2])} → {pct(ov[2])} : 최악의 손실을 시장(KOSPI {pct(mk[2])}) 수준까지 방어 (핵심 성과)", GREEN, True),
            ("수익률", f"연 {pct(bm[0])} → {pct(ov[0])} : 방어를 얻으면서 오히려 수익 상승 (현금 2.5% 운용 반영)", GREEN, True),
            ("위험대비수익", f"Sharpe {bm[1]:.2f} → {ov[1]:.2f} (시장 {mk[1]:.2f}) 로 개선", GREEN, True),
            ("vs 시장(KOSPI)", f"수익은 더 높고(vs {pct(mk[0])}), 낙폭은 시장과 비슷한 수준", BLUE, True),
            ("vol-targeting", "평시엔 리스크 스위치와 겹쳐 추가효과 작으나, 드문 꼬리위험 대비 '보험'으로 유지", GRAY, True),
            ("신뢰성", "미래 정보 미사용(실시간 재현) · 결과 재현성 확인 완료", GRAY, True)],
        Inches(0.55), Inches(1.5), Inches(12.3), Inches(5.4), size=14.5, gap=11)

# 2. 무엇을 했나 (개념 + 누적곡선)
s = slide(); header(s, "②", "무엇을 했나 — '리스크 오버레이'", "좋은 종목을 고르는 일과, 시장 위험에 따라 비중을 조절하는 일을 분리")
bullets(s, [("종목 선정", "기존 이현 스크리닝 전략 그대로 (건드리지 않음)", NAVY, True),
            ("비중 조절", "시장이 위험하면 주식 비중 ↓ (최저 20%), 안전하면 다시 ↑", NAVY, True),
            ("판단 근거", "시장 내부 건강도(상승종목 비율·신저가·추세)로 위험 국면을 읽음", NAVY, True),
            ("현금 운용", "줄인 비중만큼의 현금은 연 2.5%로 운용 (현금 관리 수익도 펀드 수익)", NAVY, True),
            ("효과", "하락장에서 덜 실어 낙폭 방어 + 상승장 수익 보존 + 현금 이자까지", GREEN, True)],
        Inches(0.55), Inches(1.55), Inches(6.2), Inches(4.8), size=13.5, gap=11)
s.shapes.add_picture(str(ASSET / "equity.png"), Inches(7.0), Inches(1.7), height=Inches(3.4))
_, tf = tb(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.7))
run(tf.paragraphs[0], "※ '이현 스크리닝 전략' = 기존 멀티팩터 스크리닝 모델. FCF_YIELD은 15개 팩터 중 하나(비중 15%)이며, 순수 FCF 유니버스가 아님.", 11.5, GRAY, italic=True)

# 3. 핵심 성과 — 낙폭 방어
s = slide(); header(s, "③", "핵심 성과 — 시장과 비교", "시장보다 수익 높고, 낙폭은 시장 수준으로 (수익·위험 둘 다 잡음)")
mt = [("전략", "연수익률", "최대낙폭", "위험대비수익", "평균 주식비중"),
      ("KOSPI 지수 (시장)", pct(mk[0]), pct(mk[2]), f"{mk[1]:.2f}", "100%"),
      ("이현 스크리닝 (단독)", pct(bm[0]), pct(bm[2]), f"{bm[1]:.2f}", "100%"),
      ("+ 리스크 오버레이 (채택)", pct(ov[0]), pct(ov[2]), f"{ov[1]:.2f}", "66%"),
      ("+ vol-targeting", pct(va[0]), pct(va[2]), f"{va[1]:.2f}", "64%")]
table(s, mt, Inches(0.55), Inches(1.55), Inches(6.7), [3.2, 1.6, 1.7, 1.8, 1.9], hl=3, row_h=0.62, fs=12)
s.shapes.add_picture(str(ASSET / "mdd.png"), Inches(7.5), Inches(1.55), height=Inches(3.1))
bullets(s, [("→", "오버레이: 수익 시장·단독 모두 상회 + 낙폭은 시장 수준(-33%)으로 방어", GREEN, True),
            ("→", "단독 전략의 약점은 '수익'이 아니라 '낙폭(-44%)' → 오버레이가 그 약점을 정확히 보완", BLUE, False),
            ("→", "남은 과제: 낙폭이 아직 -33% → slow bear(2022형) 방어를 더 개선하면 완성", GRAY, False)],
        Inches(0.55), Inches(5.0), Inches(12.3), Inches(1.9), size=12.5, gap=5)
_, tf = tb(s, Inches(0.55), Inches(6.75), Inches(12.2), Inches(0.5))
run(tf.paragraphs[0], "※ 오버레이 수익률은 미투자 현금 연 2.5% 운용 가정 반영.", 11, GRAY, italic=True)

# 4. vol-targeting
s = slide(); header(s, "④", "vol-targeting은 왜 유지하나 — '꼬리위험 보험'", "평시엔 리스크 스위치와 겹치나, 드물게 벌어지는 순간이 곧 위험관리가 필요한 때")
bullets(s, [("평시", "시장 변동성이 커지는 때 = 리스크 스위치가 이미 주식을 줄인 때 → 대부분 겹침(약 99%)", NAVY, True),
            ("드문 1%", "변동성은 튀는데 국면신호가 아직 안 잡은 '꼬리 상황' → 바로 이때가 위험관리가 필요한 순간", RED, True),
            ("결론", "평시 수익 기여는 작지만, 값싼 '보험'으로 유지 → 채택", GREEN, True)],
        Inches(0.55), Inches(1.6), Inches(6.2), Inches(3.6), size=14, gap=14)
s.shapes.add_picture(str(ASSET / "exp.png"), Inches(7.0), Inches(1.7), height=Inches(3.4))
_, tf = tb(s, Inches(0.55), Inches(6.2), Inches(12.2), Inches(0.8))
run(tf.paragraphs[0], "※ 한국시장 특성상 변동성과 국면신호가 평소 함께 움직여 겹침이 큼. 겹치지 않는 드문 구간의 대비가 유지 이유.", 12, GRAY, italic=True)

# 5. 믿을 만한가
s = slide(); header(s, "⑤", "믿을 만한가 — 점검 결과", "실전 적용 전 반드시 확인하는 세 가지")
bullets(s, [("미래 정보 미사용", "매 시점 '그때까지의 데이터'만으로 판단 → 실시간에도 동일하게 재현", NAVY, True),
            ("결과 재현성", "초기값을 바꿔도 대부분 동일한 결과. 일부 예외는 확인했고 안정된(다수) 해로 고정", NAVY, True),
            (" ", "→ 실무 적용 시 초기값 여러 개로 최적안을 고르는 방식으로 근본 보강 가능", GRAY, False),
            ("표본의 한계", "검증기간 약 8년. 낙폭 방어는 뚜렷, 수익 개선폭은 표본이 더 쌓여야 확답 가능", NAVY, True)],
        Inches(0.55), Inches(1.6), Inches(12.3), Inches(4.8), size=14.5, gap=14)

# 6. 결론 & 다음
s = slide(); header(s, "⑥", "결론 & 다음 과제", "MDD 방어는 거의 완성 — 다음은 수익률 개선 트랙")
bullets(s, [("결론", "리스크 오버레이: 시장 대비 수익↑·낙폭은 시장 수준 → 채택. vol-targeting은 꼬리보험으로 유지", GREEN, True)],
        Inches(0.55), Inches(1.55), Inches(12.3), Inches(1.0), size=15, gap=10)
ln = s.shapes.add_shape(1, Inches(0.55), Inches(2.75), Inches(12.25), Pt(1.2))
ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC); ln.line.fill.background()
bullets(s, [("남은 과제 1 · 낙폭", "slow bear(2022형 완만한 하락장) 방어 개선 → MDD 레짐모델 최종 완성", NAVY, True),
            ("남은 과제 2 · 수익 (현금)", "미투자 현금 운용(cash management) 실제화 → 약 +1%p CAGR 기여", NAVY, True),
            ("남은 과제 3 · 수익 (종목)", "모멘텀 종목 유니버스 개발 + 가치/모멘텀을 국면 따라 배분하는 AI 에이전트", NAVY, True)],
        Inches(0.55), Inches(3.0), Inches(12.3), Inches(3.2), size=14.5, gap=13)

prs.save(str(OUT))
print(f"저장 완료: {OUT}")
print(f"KOSPI {pct(mk[0])}/{mk[1]:.2f}/{pct(mk[2])} · 이현 {pct(bm[0])}/{bm[1]:.2f}/{pct(bm[2])} · 오버레이 {pct(ov[0])}/{ov[1]:.2f}/{pct(ov[2])}")
