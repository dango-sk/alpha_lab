# -*- coding: utf-8 -*-
"""레짐 모델 정리 PPT (4장) — 전문 톤 + 실험 표 + 차트. 박사님 보고용."""
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm
import numpy as np
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

for cand in ["AppleGothic", "Apple SD Gothic Neo", "NanumGothic"]:
    if cand == "AppleGothic" or any(cand in f.name for f in fm.fontManager.ttflist):
        plt.rcParams["font.family"] = cand; break
plt.rcParams["axes.unicode_minus"] = False
CH = Path("analysis/ppt_charts"); CH.mkdir(exist_ok=True)
G="#1B6E3C"; R="#B12A2A"; B="#2E5B9A"; GY="#888888"; O="#C77A14"

# C1: seed robustness
seeds=["1","2","3","4","5"]; no_fx=[-0.93,-1.75,-2.18,1.75,-1.91]; fx=[1.84,1.53,1.53,1.36,1.53]
x=np.arange(5); w=.38
fig,ax=plt.subplots(figsize=(6.6,4.0))
ax.bar(x-w/2,no_fx,w,label="내부지표만",color=R,alpha=.85); ax.bar(x+w/2,fx,w,label="+ 환율",color=G,alpha=.9)
ax.axhline(0,color="k",lw=.8); ax.set_xticks(x); ax.set_xticklabels(seeds); ax.set_xlabel("random seed (5개)")
ax.set_ylabel("Bull/Bear 수익격차 (%p)")
ax.set_title("Full-Cov HMM의 seed 민감도와 환율의 안정화 효과", fontsize=10.5)
ax.legend(fontsize=10); fig.tight_layout(); fig.savefig(CH/"p1.png",dpi=150); plt.close()

# C2: config
cfgs=["기본\n(환율)","+금리","+달러","+뉴스","3-state","4-state"]; mu=[1.56,2.32,1.77,-0.80,-0.58,0.43]; sd=[0.16,0.12,0.0,3.29,1.82,1.90]
note=["채택","FCF 무익","약효과","불안정","불안정","불안정"]; cols=[G,O,B,R,R,R]
fig,ax=plt.subplots(figsize=(6.6,4.0))
ax.bar(np.arange(6),mu,yerr=sd,color=cols,alpha=.85,capsize=4); ax.axhline(0,color="k",lw=.8)
for i,(m,sdv,t) in enumerate(zip(mu,sd,note)): ax.text(i,m+(sdv+0.25 if m>=0 else -sdv-0.55),t,ha="center",fontsize=8.5,color=cols[i])
ax.set_xticks(np.arange(6)); ax.set_xticklabels(cfgs,fontsize=9); ax.set_ylabel("수익격차 5-seed 평균±표준편차 (%p)")
ax.set_title("피처/구조 추가 실험: 기본(환율)만 안정·양(+)", fontsize=10.5)
fig.tight_layout(); fig.savefig(CH/"p2.png",dpi=150); plt.close()

# C3: vs ai_v2
fig,axs=plt.subplots(1,2,figsize=(7.2,3.8)); xx=np.arange(2); w=.36
axs[0].bar(xx-w/2,[7,1.53],w,label="HSMM+환율",color=G); axs[0].bar(xx+w/2,[6,0.23],w,label="ai_v2",color=GY)
axs[0].set_xticks(xx); axs[0].set_xticklabels(["Event Recall\n(/7)","수익격차\n(%p)"],fontsize=9); axs[0].set_title("레짐 품질",fontsize=10); axs[0].legend(fontsize=9)
axs[1].bar(xx-w/2,[387,37.2],w,label="HSMM+환율\n(overlay)",color=G); axs[1].bar(xx+w/2,[406,40.6],w,label="ai_v2",color=GY)
axs[1].set_xticks(xx); axs[1].set_xticklabels(["누적수익\n(%)","최대낙폭\n(%)"],fontsize=9); axs[1].set_title("FCF 성과 (overlay 적용)",fontsize=10); axs[1].legend(fontsize=8.5)
fig.tight_layout(); fig.savefig(CH/"p3.png",dpi=150); plt.close()

# ---- 수식 이미지 렌더 (mathtext) ----
def formula(name, tex, w=9.5, h=0.62, fs=19):
    fig=plt.figure(figsize=(w,h)); fig.patch.set_alpha(0)
    fig.text(0.005,0.5,tex,fontsize=fs,va='center',ha='left')
    fig.savefig(CH/f"f_{name}.png",dpi=200,transparent=True,bbox_inches='tight',pad_inches=0.04); plt.close()
formula("notation", r"$x_t \sim \mathcal{N}(\mu_{s_t},\,\Sigma_{s_t}), \quad A_{ij}=P(s_{t+1}{=}j \mid s_t{=}i)$")
formula("viterbi", r"$\hat{s}=\mathrm{argmax}_{s}\;\sum_t\,[\,\log A_{s_{t-1}s_t}+\log \mathcal{N}(x_t;\mu_{s_t},\Sigma_{s_t})\,]$")
formula("sticky", r"$A_{i\cdot}\sim \mathrm{Dir}(\alpha),\;\; \alpha_{ii}{=}\kappa{=}50,\;\; \alpha_{ij}{=}1\,(i{\neq}j)\;\Rightarrow\; A_{ii}\!\uparrow$")
formula("penalty", r"$\log A'_{ij}=\log A_{ij}-c\,[\,i{\neq}j\,], \qquad c=2.0$")
formula("hyst", r"$\gamma_t{=}P(s_t{=}\mathrm{Bear}\mid x_{1:t});\;\; \gamma_t{\geq}0.65{\Rightarrow}\mathrm{Bear},\; \gamma_t{<}0.35{\Rightarrow}\mathrm{Bull},\; \mathrm{else}\;\hat r_{t-1}$", w=11.5)
formula("hsmm1", r"$\mathrm{HMM}:\; P(d)=A_{ii}^{\,d-1}(1-A_{ii}) \quad (\mathrm{geometric})$")
formula("hsmm2", r"$\mathrm{HSMM}:\; p_i(d)\sim\mathrm{NegBin};\;\; \max\!\sum_{\mathrm{seg}}[\,\log p_s(d)+\!\!\sum_{t\in\mathrm{seg}}\!\!\log\mathcal{N}(x_t;\mu_s,\Sigma_s)\,]$", w=11.5)

# 검증② 차트 (Bull/Bear 실제 수익)
fig,ax=plt.subplots(figsize=(4.5,3.5))
vals=[2.00,0.46]; ax.bar(["Bull로\n분류한 달","Bear로\n분류한 달"],vals,color=[G,R],alpha=.85)
for i,v in enumerate(vals): ax.text(i,v+0.04,f"+{v:.2f}%",ha="center",fontsize=12,fontweight="bold")
ax.set_ylabel("실제 다음달 평균수익 (%)"); ax.set_ylim(0,2.5)
ax.set_title("분류의 실제 수익 (격차 +1.53%p)",fontsize=10.5)
fig.tight_layout(); fig.savefig(CH/"valid.png",dpi=150); plt.close()

# ===== PPT =====
NAVY=RGBColor(0x1F,0x33,0x55); BLUEc=RGBColor(0x2E,0x5B,0x9A); GR=RGBColor(0x40,0x40,0x40)
REDc=RGBColor(0xB1,0x2A,0x2A); GRNc=RGBColor(0x1B,0x6E,0x3C); HDR=RGBColor(0x2E,0x5B,0x9A); HL=RGBColor(0xE8,0xF1,0xE8)
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)

def slide(title,sub=None):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    tf=s.shapes.add_textbox(Inches(0.45),Inches(0.28),Inches(12.4),Inches(0.95)).text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=title; p.font.size=Pt(24); p.font.bold=True; p.font.color.rgb=NAVY
    if sub:
        q=tf.add_paragraph(); q.text=sub; q.font.size=Pt(12.5); q.font.italic=True; q.font.color.rgb=BLUEc
    return s
def body(s,items,left,top,width,size=13):
    tf=s.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(5.6)).text_frame; tf.word_wrap=True
    for i,(lvl,txt,col,bold) in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=txt; p.level=lvl; p.font.size=Pt(size-lvl*1.5); p.font.bold=bold; p.font.color.rgb=col; p.space_after=Pt(4)
def txt(s,text,left,top,width,size=12,color=GR,bold=False):
    tf=s.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(0.5)).text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(size); p.font.color.rgb=color; p.font.bold=bold
def pic(s,name,left,top,width):
    s.shapes.add_picture(str(CH/f"f_{name}.png"),Inches(left),Inches(top),width=Inches(width))
def table(s,data,left,top,width,height,hi=None,fs=11.5,cw=None):
    gf=s.shapes.add_table(len(data),len(data[0]),Inches(left),Inches(top),Inches(width),Inches(height)); t=gf.table
    if cw:
        for j,wd in enumerate(cw): t.columns[j].width=Inches(wd)
    for i,row in enumerate(data):
        for j,v in enumerate(row):
            c=t.cell(i,j); c.text=str(v); pr=c.text_frame.paragraphs[0]
            pr.font.size=Pt(fs); pr.alignment=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT
            if i==0: pr.font.bold=True; pr.font.color.rgb=RGBColor(0xFF,0xFF,0xFF); c.fill.solid(); c.fill.fore_color.rgb=HDR
            else: c.fill.solid(); c.fill.fore_color.rgb=(HL if hi==i else RGBColor(0xFF,0xFF,0xFF)); pr.font.bold=(hi==i)
    return t

# Slide 0 — 데이터 & 용어
s=slide("0. 데이터 기간 · 용어 정의 (먼저 보기)",
        "레짐 모델 구축에 쓴 데이터와, 이후 등장하는 평가지표 용어")
txt(s,"■ 데이터",0.5,1.35,6,15,NAVY,True)
body(s,[
 (0,"기간: 2000~2026 (월단위, 약 316개월)",GR,False),
 (0,"대상: KOSPI 전종목 + KOSPI 지수",GR,False),
 (0,"입력 산출: 개별종목 종가로 breadth·52주신저가",GR,False),
 (1,"(2000~16 수집분 + 2017~ DB), 지수 모멘텀, 환율(USD/KRW, 2003~)",GR,False),
 (0,"학습: walk-forward (매월 확장 재적합, 최소 36개월)",GR,False),
 (1,"→ 그 시점까지 데이터만 사용, 미래정보 누수 없음",GR,False),
 (0,"평가구간: 2018~2026 (AI v2·뉴스 가용 구간)",BLUEc,True),
], 0.7, 1.85, 6.3, 12.5)
txt(s,"■ 용어",6.95,1.35,6,15,NAVY,True)
body(s,[
 (0,"Event Recall: 실제 위기(이후 6개월 −15%↓) 중 포착 비율 (7/7=다 잡음)",GR,False),
 (0,"Lead: 위기 시작보다 며칠 먼저 신호 켰나 (+선행/−지각, 개월)",GR,False),
 (0,"False Alarm: 위기 아닌데 약세 신호 켠 횟수(헛경보)",GR,False),
 (0,"Whipsaw: 강세↔약세 전환 총 횟수 (잦으면 거래비용↑)",GR,False),
 (0,"분별력(격차): 강세 분류달 vs 약세 분류달 실제 수익차(%p)",GR,False),
 (0,"robustness: 랜덤 초기값(seed) 바꿔도 결과 안 흔들림",GR,False),
 (0,"MDD 최대낙폭 · Sharpe 위험대비수익 · Calmar 수익/MDD",GR,False),
 (0,"Risk Overlay: 약세장에 전략 교체 대신 '비중만 축소'",GR,False),
], 7.15, 1.85, 6.0, 11.5)

# Slide 1 — 모델 구조 선택
s=slide("① 모델 구조 선택 — HSMM 채택 근거",
        "Full-Cov HMM은 분별력은 우수하나 전환 빈도(Whipsaw)가 과다 → 안정화 기법 5종 동일조건 비교")
table(s,[
 ["안정화 기법","Event Recall","Lead(개월)","False Alarm","Whipsaw","Bull/Bear 격차"],
 ["기본 Full-Cov HMM","7/7","+0.0","5","35","+1.65%p"],
 ["Sticky transition","7/7","+0.4","5","23","-0.70%p (분별력 붕괴)"],
 ["Transition penalty","7/7","-1.1","4","23","+1.40%p"],
 ["Posterior hysteresis","7/7","-1.0","5","31","+1.67%p"],
 ["HSMM (duration 학습) ★","7/7","+0.0","3","31","+1.75%p (최고)"],
], 0.5, 1.75, 12.3, 2.5, hi=5, cw=[3.2,1.9,1.7,1.9,1.6,2.0])
body(s,[
 (0,"■ 채택: HSMM — Event Recall(7/7)·Lead 유지하면서 분별력 최고(+1.75%p)·False Alarm 최저(3건)",GRNc,True),
 (1,"분별력(Bull/Bear 격차) = '약세'로 분류한 달과 '강세'로 분류한 달의 익월 수익률 차이",GR,False),
 (1,"동일 HSMM이라도 duration 학습 제거 시 분별력 +1.53→+0.77%p로 반감 → duration 모델링이 핵심 기여",GR,False),
], 0.5, 4.55, 12.3)

# Slide 1-부록(1) — 표기 · Sticky · Penalty (수식)
s=slide("①-부록(1): 안정화 기법 수식 — 표기 · Sticky · Penalty",
        "공통: 2-state full-covariance Gaussian HMM + Viterbi 디코딩")
txt(s,"기호:  x=월별 피처,  s=숨은 상태(강세/약세),  A=상태 전이확률,  μ·Σ=상태 평균·공분산,  γ=약세 사후확률",0.5,1.18,12.3,10.5,GR,False)
txt(s,"■ 공통 — HMM이 상태를 추정하고, Viterbi로 '지금 어느 상태인지' 디코딩",0.5,1.5,12.3,13,NAVY,True)
pic(s,"notation",0.9,1.92,6.8); pic(s,"viterbi",0.9,2.42,7.6)
txt(s,"■ ① Sticky — '자기 상태에 머물 확률'을 prior로 강제로 키움 (κ가 클수록 끈끈)",0.5,3.2,12.3,13,NAVY,True)
pic(s,"sticky",0.9,3.6,6.8)
txt(s,"  결과: 전환 35→23 줄었으나 분별력 +1.65→ −0.70 붕괴 ❌  (너무 끈끈해 회복장까지 '약세'에 갇힘)",0.6,4.2,12.2,11.5,REDc)
txt(s,"■ ② Penalty — 상태를 '바꿀 때마다' 점수를 깎아(−c) 전환을 억제",0.5,4.9,12.3,13,NAVY,True)
pic(s,"penalty",0.9,5.3,5.0)
txt(s,"  결과: 전환 23로 줄었으나 위기 인식이 1.1개월 늦어짐·분별력 +1.40 ❌  (안정 ↔ 선행 trade-off)",0.6,5.9,12.2,11.5,REDc)

# Slide 1-부록(2) — Hysteresis · HSMM (수식)
s=slide("①-부록(2): 안정화 기법 수식 — Hysteresis · HSMM(채택)")
txt(s,"■ ③ Hysteresis — '약세 확신(γ)'이 높아야 진입, 충분히 낮아야 탈출 (불감대로 깜빡임 억제)",0.5,1.3,12.3,13,NAVY,True)
pic(s,"hyst",0.9,1.74,8.6)
txt(s,"  결과: 전환 31·분별력 +1.67(양호) 이나 위기 인식이 1.0개월 늦음 ❌",0.6,2.36,12.2,11.5,REDc)
txt(s,"■ ④ HSMM (채택) — 각 상태가 '얼마나 오래 가는지(체류기간 d)'를 데이터에서 직접 학습",0.5,3.05,12.3,13,NAVY,True)
txt(s,"  일반 HMM은 체류기간이 기하분포로 암묵 고정 → HSMM은 학습된 분포로 디코딩",0.6,3.48,12.2,11.5,GR)
pic(s,"hsmm1",0.9,3.88,6.0); pic(s,"hsmm2",0.9,4.42,9.0)
txt(s,"  결과: 분별력 +1.75(최고)·헛경보 3(최저)·위기인식 0.0개월·Recall 7/7 유지 ✅",0.6,5.12,12.2,12,GRNc,True)
txt(s,"→ 핵심: 앞 3개는 전환을 '인위적으로 막아' 부작용(분별력↓ or 선행성↓). HSMM만 '정상 체류기간 학습'으로 부작용 없이 안정화.",0.5,5.85,12.4,12.5,NAVY,True)

# Slide 2 — 팩터 결정
s=slide("② 입력 팩터 결정 — 여러 후보 중 '환율'만 채택, 최종 6개",
        "추가 후보들은 분별력·안정성·FCF 기준을 못 넘어 전량 기각")
txt(s,"■ 최종 입력 6개  (= 시장 내부 5 + 자본흐름 1)",0.5,1.5,7.0,16,NAVY,True)
body(s,[
 (0,"①  breadth (상승종목 비율)",GR,False),
 (0,"②  Δbreadth (그 변화)",GR,False),
 (0,"③  52주 신저가 비율",GR,False),
 (0,"④  Δ신저가",GR,False),
 (0,"⑤  모멘텀 (1M − 6M)",GR,False),
 (0,"⑥  환율 Δ3M   ★ 유일하게 채택된 '추가' 후보",GRNc,True),
], 0.8, 2.05, 6.6, 16)
txt(s,"추가 후보 검증 결과 (환율 외 전량 기각)",7.8,1.55,5.3,12.5,GR,True)
table(s,[
 ["추가 후보","핵심 결과","채택"],
 ["환율","seed 안정성 대폭 개선 ★","✅"],
 ["금리(US 10Y)","분별력↑이나 FCF 무익","❌"],
 ["달러(DXY)","약효과","❌"],
 ["뉴스(LLM)","robustness 파괴","❌"],
 ["구리/금","예측 무익·후행","❌"],
], 7.8, 2.0, 5.1, 2.5, hi=1, fs=11.5, cw=[1.7,2.6,0.8])
txt(s,"→ '정보를 더하면 좋아진다'가 아니었음. 환율은 예측력이 아니라 'seed 안정성(robustness)'을 줘서 채택.",0.5,5.4,12.4,13.5,NAVY,True)

# Slide 2-부록 — 변수 투입 방식
s=slide("②-부록: ②의 후보들을 '어떻게' 넣었나 (투입 방식)",
        "② = '무엇을' 넣을지 / 여기 = '어떻게' 넣었나 — 모든 후보 동일 절차, 차이는 변환·차원뿐")
body(s,[
 (0,"■ 공통 절차",NAVY,True),
 (1,"① 기존 6피처 [breadth·Δbreadth·신저가·Δ신저가·모멘텀·환율]에 후보를 '열(column)'로 추가(hstack)",GR,False),
 (1,"② 전부 '3개월 변화'로 변환 — 레벨은 추세(비정상)라 HMM emission(정상성 가정)에 부적합",GR,False),
 (1,"③ 표준화(StandardScaler)는 학습구간으로만 fit → 미래정보 누수 없음",GR,False),
 (1,"④ Full-Cov emission: 추가 열이 상태의 평균 + 공분산(타 피처와의 상관)을 함께 형성",GR,False),
 (1,"⑤ Bear 상태 식별은 항상 breadth 5개 점수로 — 추가변수는 '상태 형성·안정화'만, 이름표는 안 붙임",GR,False),
 (1,"⑥ walk-forward(매월 재적합) + 5-seed 다수결, point-in-time(직전월 뉴스만)",GR,False),
], 0.5, 1.55, 12.3, 12.5)
table(s,[
 ["후보 변수","원자료","변환 형태","투입 차원","채택"],
 ["환율","USD/KRW","3개월 변화율(%)  (+원화약세)","6 (기본 포함)","✅"],
 ["금리","미국 10Y 수익률","3개월 차이(pp)  (+금리상승)","7","❌"],
 ["달러","DXY 달러지수","3개월 변화율(%)  (+달러강세)","7","❌"],
 ["뉴스","LLM 점수","Tightening·Stress 0~100 (레벨)","8 (2열)","❌"],
 ["구리/금","구리/금 비율","3개월 변화율(%)  (−경기둔화)","검증 탈락","❌"],
], 0.5, 4.35, 12.3, 2.6, hi=1, fs=11.5, cw=[1.7,2.5,4.2,2.0,1.0])
body(s,[
 (0,"→ 즉 '3개월 변화로 만들어 → base 옆에 열로 붙여 → 표준화 → full-cov HMM emission 학습'. 차이는 차원(7/8)·변환 형태뿐.",GRNc,True),
], 0.5, 6.95, 12.3, 12.5)

# Slide 2-부록2 — 레짐 라벨 검증
s=slide("②-부록2: 'Bear=breadth로 라벨'(앞 장 ⑤) — 그 정의가 타당한가?",
        "이름표는 breadth로(미래 미사용) 붙임 → 실제 위기·수익과 '독립 대조'로 검증")
txt(s,"■ 검증①  위기 포착 — 가격으로만 정의한 '진짜 위기'를 다 잡았나",0.5,1.35,12.2,14,NAVY,True)
body(s,[
 (1,"독립 위기 정의: 그 달 이후 6개월 내 KOSPI −15% 이상 하락 (breadth와 무관, 순수 가격)",GR,False),
 (1,"2018~2026 위기 7건 → 7건 전부 포착  =  Event Recall 7/7 (100%)",GRNc,True),
 (1,"예: 2018 긴축 · 2020 코로나 · 2021~22 긴축 · 2024 · 2025 하락 모두 ✓",GR,False),
], 0.8, 1.85, 7.2, 13)
txt(s,"■ 검증②  수익 분리 — 'Bear'라 한 달이 실제로 덜 벌었나",0.5,3.75,12.2,14,NAVY,True)
body(s,[
 (1,"Bull 분류 달 실제 다음달 평균 +2.00%  vs  Bear 분류 달 +0.46%",GR,False),
 (1,"→ 격차 +1.53%p (양수): 이름표가 실제 수익 방향과 일치 (우측 그래프)",GRNc,True),
], 0.8, 4.25, 7.2, 13)
s.shapes.add_picture(str(CH/"valid.png"), Inches(8.4), Inches(2.0), width=Inches(4.4))
txt(s,"→ 라벨은 미래를 안 보고 breadth로 붙였으나 실제 위기·실제 수익과 모두 일치 → '자의적'이 아니라 '검증된' 라벨.",0.5,6.4,12.4,12.5,GRNc,True)

# Slide 3 — 환율 robustness
s=slide("③ 환율의 역할 — 예측력이 아닌 'robustness'",
        "내부지표만으로는 random seed에 따라 분별력 부호가 역전 → 환율이 이를 안정화")
body(s,[
 (0,"■ 관찰",NAVY,True),
 (1,"내부지표만: seed별 Bull/Bear 격차가 −2.18 ~ +1.75로 부호 역전 (분별력이 초기값에 좌우)",REDc,False),
 (1,"+환율: 5개 seed 전부 +1.4~+1.8%p로 수렴 (우측)",GRNc,True),
 (0,"■ 해석",NAVY,True),
 (1,"환율은 시장 내부정보에 '자본흐름'이라는 직교 축을 추가",GR,False),
 (1,"→ HMM 우도면의 다봉성을 완화, 동일 해로 수렴 = 안정화",GR,False),
 (0,"→ 환율의 채택 사유는 '분별력 향상'이 아니라 '재현성(robustness) 확보'.",NAVY,True),
],0.5,1.6,6.0)
s.shapes.add_picture(str(CH/"p1.png"), Inches(6.8), Inches(1.9), width=Inches(6.2))

# Slide 4 — 레짐 탐지 시계열 (빨강/파랑)
s=slide("④ 레짐 탐지 결과 — 시계열 (빨강=약세 / 파랑=강세)",
        "KOSPI(log) 위에 walk-forward 레짐 표시 · 2008·2011·2018·2020·2022 위기 구간 포착")
s.shapes.add_picture(str(CH/"p_timeline.png"), Inches(0.55), Inches(1.6), width=Inches(12.2))
body(s,[
 (0,"빨강(약세) 구간이 2008 GFC·2011 유럽·2018 긴축·2020 코로나·2022 긴축 등 주요 하락 국면과 정합.",GR,False),
 (0,"환율 결합으로 seed에 무관하게 동일한 시계열 산출(재현성 확보).",GRNc,True),
], 0.55, 6.5, 12.2, 12.5)

# Slide 5 — 결과 + overlay
s=slide("⑤ 결과 및 개선 — ai_v2 대비 우위/열위와 Risk Overlay",
        "overlap 2018~2026")
body(s,[
 (0,"■ vs ai_v2: 레짐 품질(Event Recall 7/7 vs 6/7, 분별력 +1.53 vs +0.23)·낙폭방어 우위 / 단순 적용 시 절대수익은 열위",GR,False),
 (1,"원인(월별 attribution): ① 2022 긴축장에서 내부지표가 반등에 강세로 오인(−7.4%p) ② 회복장 손절이 저점 매도(−5.2%p, 139건)",GR,False),
 (0,"■ 개선: 약세 신호 시 전략 전면교체 대신 강세전략 유지 + 익스포저 축소(Risk Overlay)",NAVY,True),
],0.5,1.2,12.5)
table(s,[
 ["운용 방식","누적수익","Sharpe","최대낙폭(MDD)"],
 ["강세전략 단독","368%","0.78","-43.7%"],
 ["레짐 전략교체 (기존)","353%","0.75","-38.1%"],
 ["Risk Overlay (Bear 익스포저 50%) ★","387%","0.82","-37.2%"],
 ["ai_v2 (뉴스 기반)","406%","0.82","-40.6%"],
], 0.5, 2.85, 8.4, 2.2, hi=3, fs=12, cw=[4.0,1.6,1.3,1.5])
body(s,[
 (0,"→ Risk Overlay: 강세전략 단독 대비 수익↑·MDD↓, ai_v2와 Sharpe 동률·MDD는 우위.",GRNc,True),
 (1,"ai_v2의 높은 수익은 뉴스 '정보' 효과이며, 본 모델은 그 정보 없이 익스포저 조절만으로 근접 (거래비용 반영치).",GR,False),
],0.5,5.25,12.3)
s.shapes.add_picture(str(CH/"p3.png"), Inches(9.0), Inches(2.7), width=Inches(4.0))

out=Path("레짐모델_정리.pptx"); prs.save(out)
print(f"저장: {out.resolve()} (4장)")
