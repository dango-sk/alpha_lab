# -*- coding: utf-8 -*-
"""박사님 논의 질문 — 별도 PPT (2장): A. 모델 개선 과제 / B. 방법론 검증."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

NAVY=RGBColor(0x1F,0x33,0x55); BLUEc=RGBColor(0x2E,0x5B,0x9A); GR=RGBColor(0x40,0x40,0x40)
GRNc=RGBColor(0x1B,0x6E,0x3C); STAR=RGBColor(0xB1,0x2A,0x2A)
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)

def slide(title,sub):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    tf=s.shapes.add_textbox(Inches(0.5),Inches(0.3),Inches(12.3),Inches(0.95)).text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=title; p.font.size=Pt(25); p.font.bold=True; p.font.color.rgb=NAVY
    q=tf.add_paragraph(); q.text=sub; q.font.size=Pt(12.5); q.font.italic=True; q.font.color.rgb=BLUEc
    return s
def body(s,items,left,top,width):
    tf=s.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(5.8)).text_frame; tf.word_wrap=True
    for i,(txt,col,bold,sz) in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=txt; p.font.size=Pt(sz); p.font.bold=bold; p.font.color.rgb=col; p.space_after=Pt(3)

Q="질문"; C="→ 배경(우리가 부딪힌 지점)"

# Slide A
s=slide("박사님 논의 질문 (A) — HSMM 모델 개선 과제","앞으로 더 개선해야 할 기술적 과제")
body(s,[
 ("1. 분별력은 좋아졌는데 PnL로 안 이어지는 문제, 어떻게 메우나? ★",STAR,True,14),
 ("   금리·DXY로 Bull/Bear 수익격차를 +1.5→+2.4%p 올렸으나 FCF 수익은 그대로. '레짐 분별력 ↑ ≠ 전략 수익 ↑'을 잇는 정석 방법?",GR,False,11.5),
 ("2. 매크로 '정보'를 HMM에 어떻게 결합하나? ★",STAR,True,14),
 ("   ai_v2 우위는 뉴스/매크로 정보. feature로 직접 넣으니 robustness 붕괴. 정보를 emission이 아니라 레짐 게이트/오버레이로 결합하는 게 맞나?",GR,False,11.5),
 ("3. Slow Bear(2022)의 whipsaw를 어떻게 잡나?",NAVY,True,14),
 ("   내부지표(breadth)가 완만한 하락장의 반등에 '강세'로 깜빡임. duration 학습만으론 부족 — 비대칭 지속 등 처방?",GR,False,11.5),
 ("4. full-cov + 6피처를 36개월 창에 적합 — 과적합 통제 충분한가?",NAVY,True,14),
 ("   seed 취약성이 작은 표본·고차원 탓 같음. 공분산 shrinkage(Ledoit-Wolf 등)/정규화가 필요한가?",GR,False,11.5),
 ("5. 익스포저 사이징을 연속적으로?",NAVY,True,14),
 ("   현재 Bear=비중 50% 이분법. posterior 확신도에 비례한 연속 익스포저가 나은가? 최적 비중 추정법?",GR,False,11.5),
 ("6. HSMM duration 분포 명세",NAVY,True,14),
 ("   체류기간을 음이항(MoM)으로 추정. 이 분포 선택이 타당한지·더 나은 HSMM specification?",GR,False,11.5),
], 0.55, 1.5, 12.3)

# Slide B
s=slide("박사님 논의 질문 (B) — 방법론 검증·확인","만들면서 계속 헷갈렸던, 확인받고 싶은 점")
body(s,[
 ("7. '무엇을 예측해야 하나' — 방향 vs 위험국면",NAVY,True,14),
 ("   다음달 방향은 ~50%로 불가 → 목표를 '위험 국면 탐지'로 재정의. 이 재정의·Event Recall 중심 평가가 타당한가?",GR,False,11.5),
 ("8. Bear를 'breadth(시장구조)'로 정의 vs '실제 저수익'으로 정의",NAVY,True,14),
 ("   breadth 라벨링(사후검증 통과) vs forward-return 라벨링(다상태서 불안정). 레짐 ground truth를 무엇으로?",GR,False,11.5),
 ("9. 환율의 역할이 '예측력'이 아니라 'robustness'인 것 — 정상 해석인가? ★",STAR,True,14),
 ("   환율 단독 예측력은 약한데 넣으면 seed에 안 흔들림. '직교 축이 식별성(identifiability)을 높인다'는 설명이 타당한가?",GR,False,11.5),
 ("10. seed 민감도를 'multi-seed 다수결'로 푸는 게 정공법인가? ★",STAR,True,14),
 ("   EM 초기값 민감 → consensus로 안정화. 정석인지, 아니면 모델 명세를 바꿔야 하는지?",GR,False,11.5),
 ("11. 레짐 정의 호라이즌(이후 6개월·−15%)이 월단위 전략에 맞나?",NAVY,True,14),
 ("   위기를 'fwd 6m, −15% 낙폭'으로 정의. 월단위 팩터전환 운용 호라이즌과 정합하는가?",GR,False,11.5),
 ("12. overlay/거래비용을 분석적 근사로 둔 한계",NAVY,True,14),
 ("   전환마다 (1−e)×슬리피지로 근사. 백테스트 엔진에 익스포저 조절을 직접 넣어야 정확한가?",GR,False,11.5),
], 0.55, 1.5, 12.3)
# 핵심 강조
tf=prs.slides[-1].shapes.add_textbox(Inches(0.55),Inches(6.85),Inches(12.3),Inches(0.5)).text_frame
p=tf.paragraphs[0]; p.text="★ 가장 핵심: 1(분별력≠PnL) · 2(정보 결합) · 9·10(robustness 해석/해법)"; p.font.size=Pt(12.5); p.font.bold=True; p.font.color.rgb=GRNc

out=Path("박사님_논의질문.pptx"); prs.save(out)
print(f"저장: {out.resolve()} (2장)")
